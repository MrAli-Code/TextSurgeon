#!/usr/bin/env python3
"""Surgeon Agent Engine — Multi-file Code Generation, Workspace Execution & Verification.

Part of Text Surgeon. Extends precision surgical editing into a full autonomous
Agent Mode:
1. Translates high-level goals into Agent Protocol prompts for any LLM.
2. Ingests LLM replies containing multi-file creations (@@FILE), surgical edits
   (@@EDIT), dependency setups (@@COMMAND), and execution runs (@@RUN).
3. Safely writes files inside designated project workspaces with atomic backups.
4. Executes dependencies and scripts inside the workspace with live logging.
5. Emits post-execution verification and diagnostic prompts for iterative debugging.

Zero third-party dependencies: uses standard library modules only.
"""

from __future__ import annotations

import difflib
import hashlib
import io
import json
import os
import re
import shutil
import subprocess
import sys
import threading
import time
import urllib.error
import urllib.parse
import urllib.request
import uuid
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Callable, Dict, Final, List, Optional, Sequence, Set, Tuple, Union


import surgeon_engine as engine
from surgeon_engine import EditOp, SelectionError, apply_edit_ops

AGENT_PROTOCOL_VERSION: Final[str] = "1.0"
MAX_BODY_BYTES: Final[int] = 64 * 1024 * 1024
DEFAULT_EXECUTION_TIMEOUT: Final[int] = 180  # seconds


# --------------------------------------------------------------------------- #
# Exceptions
# --------------------------------------------------------------------------- #


class AgentError(Exception):
    """Base error for all Agent operations."""

    def __init__(self, message: str, hint: Optional[str] = None) -> None:
        super().__init__(message)
        self.hint = hint


class WorkspaceSecurityError(AgentError):
    """Raised when a file path attempts to escape the designated workspace."""


class AgentParseError(AgentError):
    """Raised when an AI agent response is malformed."""


# --------------------------------------------------------------------------- #
# Data Models
# --------------------------------------------------------------------------- #


@dataclass
class FileAction:
    """Action to perform on a file within the workspace.

    action_type:
      - 'create': create or overwrite full file
      - 'edit': apply surgical EditOps to existing file
      - 'delete': remove file
    """

    path: str
    action_type: str  # 'create', 'edit', 'delete'
    content: str = ""
    edit_ops: List[EditOp] = field(default_factory=list)
    strategy: Optional[str] = None


@dataclass
class AgentPlan:
    """The parsed plan of actions emitted by an LLM."""

    explanation: Optional[str]
    file_actions: List[FileAction] = field(default_factory=list)
    setup_commands: List[str] = field(default_factory=list)
    run_command: Optional[str] = None
    expected_artifacts: List[str] = field(default_factory=list)
    raw_response: str = ""


@dataclass
class ExecutionResult:
    """The outcome of running commands in the project workspace."""

    command: str
    exit_code: int
    stdout: str
    stderr: str
    duration_sec: float
    artifacts_found: List[str] = field(default_factory=list)
    files_created_or_modified: List[str] = field(default_factory=list)
    timed_out: bool = False
    success: bool = True

    def summary(self) -> str:
        """Returns a concise summary line."""
        status = "SUCCESS" if self.success else f"FAILED (exit code {self.exit_code})"
        return f"Command '{self.command}' finished in {self.duration_sec:.2f}s: {status}"


# --------------------------------------------------------------------------- #
# Workspace Security & Filesystem Helpers
# --------------------------------------------------------------------------- #


def resolve_safe_workspace_path(workspace_root: str, relative_path: str) -> str:
    """Resolves a relative path within the workspace root, preventing directory traversal.

    Args:
        workspace_root: Absolute path to the designated project directory.
        relative_path: User/AI supplied relative file path (e.g. 'src/main.py').

    Returns:
        The safe absolute path inside workspace_root.

    Raises:
        WorkspaceSecurityError: If path tries to escape workspace_root.
    """
    clean_root = os.path.abspath(workspace_root)
    norm_rel = relative_path.replace("\\", "/").strip()

    # Block raw absolute paths or drive letters unless they resolve strictly inside clean_root
    if norm_rel.startswith("/") or re.match(r"^[a-zA-Z]:", norm_rel) or os.path.isabs(relative_path):
        target = os.path.abspath(relative_path)
        try:
            common = os.path.commonpath([clean_root, target])
            if common == clean_root:
                return target
        except ValueError:
            pass
        raise WorkspaceSecurityError(
            f"Path '{relative_path}' escapes the workspace directory '{workspace_root}'.",
            hint="All files must reside inside the designated project directory.",
        )

    target = os.path.abspath(os.path.join(clean_root, norm_rel))
    try:
        common = os.path.commonpath([clean_root, target])
        if common != clean_root:
            raise WorkspaceSecurityError(
                f"Path '{relative_path}' escapes the workspace directory '{workspace_root}'.",
                hint="All files must reside inside the designated project directory.",
            )
    except ValueError:
        raise WorkspaceSecurityError(
            f"Path '{relative_path}' escapes the workspace directory '{workspace_root}'.",
            hint="All files must reside inside the designated project directory.",
        )
    return target


def scan_workspace_tree(
    workspace_root: str, max_files: int = 150, max_depth: int = 4
) -> Dict[str, Any]:
    """Scans the project directory and returns structural metadata.

    Args:
        workspace_root: Absolute path to project folder.
        max_files: Max files to inspect.
        max_depth: Max directory depth.

    Returns:
        Dict with file list, total sizes, total lines, est_tokens, detected languages, and recent artifacts.
    """
    clean_root = os.path.abspath(workspace_root)
    if not os.path.isdir(clean_root):
        return {
            "exists": False,
            "root": clean_root,
            "files": [],
            "total_files": 0,
            "total_size": 0,
            "total_lines": 0,
            "est_tokens": 0,
            "languages": [],
        }

    files_list: List[Dict[str, Any]] = []
    total_size = 0
    total_lines = 0
    languages: Set[str] = set()
    root_depth = clean_root.rstrip("\\/").count(os.sep)

    ignore_dirs = {
        ".git",
        "__pycache__",
        "node_modules",
        ".venv",
        "venv",
        "dist",
        "build",
        ".gemini",
        ".idea",
        ".vscode",
        ".surgeon",
    }

    for root, dirs, files in os.walk(clean_root):
        dirs[:] = [d for d in dirs if d not in ignore_dirs]
        current_depth = root.count(os.sep) - root_depth
        if current_depth > max_depth:
            dirs[:] = []
            continue

        for fname in sorted(files):
            if fname.startswith(".surgeon") or fname.endswith(".bak") or fname.endswith(".tmp"):
                continue
            fpath = os.path.join(root, fname)
            rel_path = os.path.relpath(fpath, clean_root).replace("\\", "/")
            try:
                stat = os.stat(fpath)
                sz = stat.st_size
                mtime = stat.st_mtime
            except OSError:
                continue

            file_lines = 0
            ext = os.path.splitext(fname)[1].lower()
            if ext in (".py", ".pyw"):
                languages.add("python")
            elif ext in (".js", ".mjs", ".cjs", ".ts"):
                languages.add("javascript/typescript")
            elif ext in (".html", ".htm"):
                languages.add("html")
            elif ext in (".css", ".scss", ".less"):
                languages.add("css")
            elif ext in (".json", ".yaml", ".yml", ".toml"):
                languages.add("config")
            elif ext in (".sh", ".bat", ".ps1"):
                languages.add("shell")
            elif ext in (".md", ".txt", ".rst"):
                languages.add("docs")
            elif ext in (".pptx", ".docx", ".xlsx", ".pdf", ".png", ".jpg", ".svg", ".csv"):
                languages.add("artifacts")

            # Try counting lines for text files
            if ext in (".py", ".pyw", ".js", ".ts", ".html", ".css", ".json", ".yaml", ".yml", ".toml", ".sh", ".md", ".txt", ".csv"):
                try:
                    with open(fpath, "rb") as fh:
                        file_lines = sum(1 for _ in fh)
                        total_lines += file_lines
                except OSError:
                    pass

            total_size += sz
            files_list.append(
                {
                    "path": rel_path,
                    "name": fname,
                    "size": sz,
                    "lines": file_lines,
                    "mtime": mtime,
                    "ext": ext,
                }
            )
            if len(files_list) >= max_files:
                break
        if len(files_list) >= max_files:
            break

    return {
        "exists": True,
        "root": clean_root,
        "files": files_list,
        "total_files": len(files_list),
        "total_size": total_size,
        "total_lines": total_lines,
        "est_tokens": max(1, total_size // 4),
        "languages": sorted(list(languages)),
    }


# --------------------------------------------------------------------------- #
# Agent Response Parser
# --------------------------------------------------------------------------- #

_EXPLANATION_RE: Final["re.Pattern[str]"] = re.compile(
    r"<EXPLANATION>(.*?)</EXPLANATION>", re.IGNORECASE | re.DOTALL
)
_COMMAND_LINE_RE: Final["re.Pattern[str]"] = re.compile(
    r"(?m)^[ \t]*@@(?:COMMAND|SETUP)[ \t]+([^\n]+)$"
)
_RUN_LINE_RE: Final["re.Pattern[str]"] = re.compile(
    r"(?m)^[ \t]*@@RUN[ \t]+([^\n]+)$"
)
_ARTIFACT_LINE_RE: Final["re.Pattern[str]"] = re.compile(
    r"(?m)^[ \t]*@@ARTIFACT[ \t]+([^\n]+)$"
)
_DELETE_LINE_RE: Final["re.Pattern[str]"] = re.compile(
    r"(?m)^[ \t]*@@DELETE[ \t]+([^\n]+)$"
)


def _clean_token(s: str) -> str:
    """Removes invisible/bidi characters from header tokens."""
    return "".join(ch for ch in s if ch not in engine._STRUCTURAL_IGNORE).strip()


def parse_agent_response(raw_text: str) -> AgentPlan:
    """Parses an LLM response containing Agent Protocol blocks or JSON.

    Supports:
    - <EXPLANATION>...</EXPLANATION>
    - @@COMMAND <cmd> / @@SETUP <cmd>
    - @@RUN <cmd>
    - @@ARTIFACT <file_pattern>
    - @@DELETE <rel_path>
    - @@FILE <rel_path>
      <<<
      file contents
      >>>
    - @@EDIT <rel_path> [strategy]
      START-ANCHOR: ...
      END-ANCHOR: ...
      <<<
      replacement
      >>>
    - Complete JSON fallback:
      { "explanation": "...", "files": [...], "commands": [...], "run": "..." }

    Args:
        raw_text: The full response string pasted from LLM.

    Returns:
        An AgentPlan instance.

    Raises:
        AgentParseError: If response is empty or unparseable.
    """
    if not raw_text or not raw_text.strip():
        raise AgentParseError("The response is empty.", hint="Paste the full reply from the AI.")

    explanation: Optional[str] = None
    remainder = raw_text

    expl_match = _EXPLANATION_RE.search(raw_text)
    if expl_match:
        explanation = expl_match.group(1).strip()
        remainder = raw_text[: expl_match.start()] + "\n" + raw_text[expl_match.end() :]

    # Check for Markdown Protocol headers
    has_file = bool(re.search(r"(?m)^[ \t]*@@FILE\b", remainder))
    has_edit = bool(re.search(r"(?m)^[ \t]*@@EDIT\b", remainder))
    has_cmd = bool(re.search(r"(?m)^[ \t]*@@(?:COMMAND|SETUP|RUN|DELETE|ARTIFACT)\b", remainder))

    if has_file or has_edit or has_cmd:
        return _parse_markdown_agent_protocol(remainder, explanation, raw_text)

    # If no @@ markers, try JSON parser
    try:
        return _parse_json_agent_protocol(remainder, explanation, raw_text)
    except AgentParseError:
        # Fallback: check if entire response is a single code fence block
        code_fence_match = re.search(r"```(?:python|javascript|js|html|css|sh|bash)?\s*\n(.*?)\n```", raw_text, re.DOTALL)
        if code_fence_match:
            # Create a single main file from the code fence
            code_body = code_fence_match.group(1)
            guessed_name = "main.py" if ("def " in code_body or "import " in code_body) else "script.js"
            return AgentPlan(
                explanation=explanation or "Auto-extracted single script from code block.",
                file_actions=[FileAction(path=guessed_name, action_type="create", content=code_body)],
                run_command=f"python {guessed_name}" if guessed_name.endswith(".py") else f"node {guessed_name}",
                raw_response=raw_text,
            )
        raise AgentParseError(
            "No @@FILE, @@EDIT, @@COMMAND or JSON payload found in the response.",
            hint="Ask the AI to use @@FILE and @@RUN blocks as specified in the Agent Protocol.",
        )


def _parse_markdown_agent_protocol(
    text: str, explanation: Optional[str], raw_full: str
) -> AgentPlan:
    """Parses markdown @@ blocks into an AgentPlan."""
    setup_commands: List[str] = []
    run_command: Optional[str] = None
    expected_artifacts: List[str] = []
    file_actions: List[FileAction] = []

    # Extract command lines (both inline and block format)
    lines = text.splitlines(keepends=False)
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i]
        clean_l = _clean_token(line)

        # @@RUN [cmd]
        run_m = re.match(r"^@@RUN(?:\s+(.+))?$", clean_l, re.IGNORECASE)
        if run_m:
            inline_cmd = (run_m.group(1) or "").strip()
            if inline_cmd:
                run_command = inline_cmd
                i += 1
                continue
            i += 1
            if i < n and _clean_token(lines[i]) == "<<<":
                i += 1
                cmd_lines: List[str] = []
                while i < n and _clean_token(lines[i]) != ">>>":
                    cmd_lines.append(lines[i])
                    i += 1
                i += 1
                run_command = "\n".join(cmd_lines).strip()
            elif i < n and lines[i].strip() and not lines[i].strip().startswith("@@"):
                run_command = lines[i].strip()
                i += 1
            continue

        # @@COMMAND / @@SETUP [cmd]
        cmd_m = re.match(r"^@@(?:COMMAND|SETUP)(?:\s+(.+))?$", clean_l, re.IGNORECASE)
        if cmd_m:
            inline_cmd = (cmd_m.group(1) or "").strip()
            if inline_cmd:
                if inline_cmd not in setup_commands:
                    setup_commands.append(inline_cmd)
                i += 1
                continue
            i += 1
            if i < n and _clean_token(lines[i]) == "<<<":
                i += 1
                cmd_lines = []
                while i < n and _clean_token(lines[i]) != ">>>":
                    cmd_lines.append(lines[i])
                    i += 1
                i += 1
                c_val = "\n".join(cmd_lines).strip()
                if c_val and c_val not in setup_commands:
                    setup_commands.append(c_val)
            elif i < n and lines[i].strip() and not lines[i].strip().startswith("@@"):
                c_val = lines[i].strip()
                if c_val not in setup_commands:
                    setup_commands.append(c_val)
                i += 1
            continue

        # @@ARTIFACT [path]
        art_m = re.match(r"^@@ARTIFACT\s+(.+)$", clean_l, re.IGNORECASE)
        if art_m:
            art = art_m.group(1).strip()
            if art and art not in expected_artifacts:
                expected_artifacts.append(art)
            i += 1
            continue

        # @@DELETE [path]
        del_m = re.match(r"^@@DELETE\s+(.+)$", clean_l, re.IGNORECASE)
        if del_m:
            del_path = del_m.group(1).strip()
            if del_path:
                file_actions.append(FileAction(path=del_path, action_type="delete"))
            i += 1
            continue
        line = lines[i]
        clean_l = _clean_token(line)

        # @@FILE <path>
        file_match = re.match(r"^@@FILE\s+(.+)$", clean_l, re.IGNORECASE)
        if file_match:
            raw_arg = file_match.group(1).strip()
            del_check = re.match(r"^([^\s]+)\s+(?:action\s*:\s*delete|\[delete\]|delete)$", raw_arg, re.IGNORECASE)
            if del_check:
                rel_path = del_check.group(1).strip("\"'")
                file_actions.append(FileAction(path=rel_path, action_type="delete"))
                i += 1
                continue
            rel_path = raw_arg.strip("\"'")
            i += 1
            # Look for <<< ... >>>
            while i < n and _clean_token(lines[i]) != "<<<":
                if lines[i].strip().startswith("@@"):
                    raise AgentParseError(
                        f"Missing opening delimiter '<<<' for @@FILE {rel_path}.",
                        hint="Place '<<<' on its own line before file content.",
                    )
                i += 1
            if i >= n:
                raise AgentParseError(
                    f"Missing opening delimiter '<<<' for @@FILE {rel_path}.",
                    hint="Place '<<<' on its own line before file content.",
                )
            i += 1  # pass <<<
            content_lines: List[str] = []
            found_close = False
            while i < n:
                if _clean_token(lines[i]) == ">>>":
                    found_close = True
                    break
                content_lines.append(lines[i])
                i += 1
            if not found_close:
                raise AgentParseError(
                    f"Missing closing delimiter '>>>' for @@FILE {rel_path}.",
                    hint="Place '>>>' on its own line after file content.",
                )
            i += 1  # pass >>>
            content = "\n".join(content_lines)
            if content_lines:
                content += "\n"
            file_actions.append(
                FileAction(path=rel_path, action_type="create", content=content)
            )
            continue

        # @@EDIT <path> [strategy]
        edit_match = re.match(r"^@@EDIT\s+([^\s]+)(?:\s+([a-zA-Z]+))?", clean_l, re.IGNORECASE)
        if edit_match:
            rel_path = edit_match.group(1).strip().strip("\"'")
            strategy_hint = (edit_match.group(2) or "anchor").strip().lower()
            i += 1
            headers: Dict[str, str] = {}
            before_lines: List[str] = []
            after_lines: List[str] = []

            while i < n:
                h_line = lines[i]
                clean_h = _clean_token(h_line)
                if clean_h == "<<<":
                    break
                if ":" in h_line:
                    k, v = h_line.split(":", 1)
                    k_clean = k.strip().upper().replace("-", "_")
                    v_clean = v.strip()
                    if k_clean == "BEFORE":
                        before_lines.append(v.lstrip())
                    elif k_clean == "AFTER":
                        after_lines.append(v.lstrip())
                    else:
                        headers[k_clean] = v_clean
                i += 1

            if i >= n:
                break
            i += 1  # pass <<<
            replace_lines: List[str] = []
            while i < n and _clean_token(lines[i]) != ">>>":
                replace_lines.append(lines[i])
                i += 1
            i += 1  # pass >>>

            replace_text = "\n".join(replace_lines)
            if replace_lines:
                replace_text += "\n"

            payload: Dict[str, Any] = {"strategy": strategy_hint, "replace": replace_text}
            if "START_ANCHOR" in headers:
                payload["start_anchor"] = headers["START_ANCHOR"]
            if "END_ANCHOR" in headers:
                payload["end_anchor"] = headers["END_ANCHOR"]
            if "SEARCH" in headers:
                payload["search"] = headers["SEARCH"]
            if before_lines:
                payload["before"] = before_lines
            if after_lines:
                payload["after"] = after_lines
            if "NAME" in headers:
                payload["name"] = headers["NAME"]
            if "TARGET_HINT" in headers:
                payload["target_hint"] = headers["TARGET_HINT"]

            try:
                op = EditOp.from_payload_item(payload, len(file_actions) + 1)
                existing_edit = next(
                    (
                        fa
                        for fa in file_actions
                        if fa.path == rel_path and fa.action_type == "edit"
                    ),
                    None,
                )
                if existing_edit:
                    existing_edit.edit_ops.append(op)
                else:
                    file_actions.append(
                        FileAction(
                            path=rel_path,
                            action_type="edit",
                            edit_ops=[op],
                            strategy=strategy_hint,
                        )
                    )
            except Exception as exc:
                raise AgentParseError(
                    f"Invalid @@EDIT block for '{rel_path}': {exc}",
                    hint="Check anchor/context syntax.",
                ) from exc
            continue

        i += 1

    return AgentPlan(
        explanation=explanation,
        file_actions=file_actions,
        setup_commands=setup_commands,
        run_command=run_command,
        expected_artifacts=expected_artifacts,
        raw_response=raw_full,
    )


def _parse_json_agent_protocol(
    text: str, explanation: Optional[str], raw_full: str
) -> AgentPlan:
    """Extracts and decodes JSON agent payload."""
    cleaned = re.sub(r"^[ \t]*```[^\n]*$", "", text, flags=re.MULTILINE).strip()
    data: Optional[Dict[str, Any]] = None

    decoder = json.JSONDecoder(strict=False)
    for match in re.finditer(r"\{", cleaned):
        try:
            val, _ = decoder.raw_decode(cleaned, match.start())
            if isinstance(val, dict) and ("files" in val or "file_actions" in val or "run" in val or "commands" in val):
                data = val
                break
        except json.JSONDecodeError:
            continue

    if not data:
        raise AgentParseError("Could not find a valid Agent JSON payload.")

    expl = data.get("explanation") or explanation
    setup_cmds = data.get("commands") or data.get("setup_commands") or []
    if isinstance(setup_cmds, str):
        setup_cmds = [setup_cmds]
    run_cmd = data.get("run") or data.get("run_command")
    artifacts = data.get("artifacts") or data.get("expected_artifacts") or []
    if isinstance(artifacts, str):
        artifacts = [artifacts]

    file_actions: List[FileAction] = []
    raw_files = data.get("files") or data.get("file_actions") or []
    for idx, item in enumerate(raw_files, start=1):
        if not isinstance(item, dict):
            continue
        path = str(item.get("path") or item.get("file") or f"file_{idx}.py").strip()
        act_type = str(item.get("action") or item.get("type") or "create").strip().lower()
        content = str(item.get("content") or item.get("code") or item.get("replace") or "")

        if act_type == "delete":
            file_actions.append(FileAction(path=path, action_type="delete"))
        elif act_type == "edit":
            edits_data = item.get("edits") or [item]
            ops: List[EditOp] = []
            for edit_idx, e_item in enumerate(edits_data, start=1):
                ops.append(EditOp.from_payload_item(e_item, edit_idx))
            file_actions.append(FileAction(path=path, action_type="edit", edit_ops=ops))
        else:
            file_actions.append(FileAction(path=path, action_type="create", content=content))

    return AgentPlan(
        explanation=expl,
        file_actions=file_actions,
        setup_commands=[str(c) for c in setup_cmds],
        run_command=str(run_cmd) if run_cmd else None,
        expected_artifacts=[str(a) for a in artifacts],
        raw_response=raw_full,
    )


# --------------------------------------------------------------------------- #
# Workspace File Manager & Execution
# --------------------------------------------------------------------------- #


SURGEON_DIR_NAME: Final[str] = ".surgeon"
DEFAULT_PROJECTS_DIR: Final[str] = "projects"


class ProjectManager:
    """Manages the central projects directory and individual safe project workspaces."""

    @staticmethod
    def get_projects_root(base_dir: Optional[str] = None) -> str:
        """Returns the absolute path to the central projects root folder."""
        if base_dir:
            root = os.path.abspath(base_dir)
        else:
            # Place 'projects' in current working directory or adjacent to repo root
            root = os.path.abspath(os.path.join(os.getcwd(), DEFAULT_PROJECTS_DIR))
        os.makedirs(root, exist_ok=True)
        return root

    @classmethod
    def list_projects(cls, base_dir: Optional[str] = None) -> List[Dict[str, Any]]:
        """Lists all safe projects found in the projects library."""
        root = cls.get_projects_root(base_dir)
        projects: List[Dict[str, Any]] = []

        try:
            entries = sorted(os.listdir(root))
        except OSError:
            return []

        for entry in entries:
            proj_path = os.path.join(root, entry)
            if not os.path.isdir(proj_path) or entry.startswith("."):
                continue

            ws = AgentWorkspace(proj_path)
            meta = ws.get_project_meta()
            tree = ws.scan()

            try:
                mtime = os.path.getmtime(proj_path)
            except OSError:
                mtime = time.time()

            projects.append(
                {
                    "name": entry,
                    "path": proj_path,
                    "created_at": meta.get("created_at") or datetime.fromtimestamp(mtime).isoformat(),
                    "last_modified": meta.get("last_modified") or datetime.fromtimestamp(mtime).isoformat(),
                    "runtime": meta.get("runtime", "python"),
                    "description": meta.get("description", ""),
                    "total_files": tree.get("total_files", 0),
                    "total_size": tree.get("total_size", 0),
                    "languages": tree.get("languages", []),
                    "backups_count": len(ws.list_backups()),
                }
            )

        return projects

    @classmethod
    def create_project(
        cls,
        name: str,
        base_dir: Optional[str] = None,
        runtime: str = "python",
        description: str = "",
    ) -> Dict[str, Any]:
        """Creates an organized safe project workspace."""
        clean_name = re.sub(r"[^\w\-\.]+", "_", name.strip()).strip("._")
        if not clean_name:
            clean_name = f"project_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

        root = cls.get_projects_root(base_dir)
        proj_path = os.path.join(root, clean_name)
        os.makedirs(proj_path, exist_ok=True)

        ws = AgentWorkspace(proj_path)
        now_iso = datetime.now().isoformat()
        meta = {
            "name": clean_name,
            "created_at": now_iso,
            "last_modified": now_iso,
            "runtime": runtime,
            "description": description,
            "goals": [],
        }
        ws.save_project_meta(meta)

        # Create starter files based on template/runtime if directory is empty
        runtime_lower = runtime.lower().strip()
        if not any(f for f in os.listdir(proj_path) if f != SURGEON_DIR_NAME):
            if runtime_lower in ("powerpoint", "ppt", "presentation"):
                meta["runtime"] = "powerpoint"
                starter_py = os.path.join(proj_path, "make_presentation.py")
                with open(starter_py, "w", encoding="utf-8") as fh:
                    fh.write(
                        f'"""\nPowerPoint Generator Agent Starter — {clean_name}\n"""\n\n'
                        "import os\nimport sys\n\n"
                        "def create_presentation():\n"
                        "    print('Building presentation for {clean_name}...')\n"
                        "    try:\n"
                        "        from pptx import Presentation\n"
                        "        prs = Presentation()\n"
                        "        title_slide = prs.slides.add_slide(prs.slide_layouts[0])\n"
                        "        title_slide.shapes.title.text = '{clean_name}'\n"
                        "        title_slide.placeholders[1].text = 'Created with Text Surgeon Agent'\n"
                        "        out_file = 'presentation.pptx'\n"
                        "        prs.save(out_file)\n"
                        "        print(f'Successfully generated: {{out_file}}')\n"
                        "    except ImportError:\n"
                        "        print('python-pptx not installed. Install via: pip install python-pptx')\n"
                        "        with open('presentation_notes.txt', 'w', encoding='utf-8') as f:\n"
                        "            f.write('# Presentation: {clean_name}\\nSlide 1: Title\\nSlide 2: Overview\\n')\n"
                        "        print('Saved fallback presentation notes.')\n\n"
                        "if __name__ == '__main__':\n"
                        "    create_presentation()\n"
                    )
                with open(os.path.join(proj_path, "requirements.txt"), "w", encoding="utf-8") as fh:
                    fh.write("python-pptx>=0.6.21\n")

            elif runtime_lower in ("web_scraper", "scraper"):
                meta["runtime"] = "web_scraper"
                starter_py = os.path.join(proj_path, "scraper.py")
                with open(starter_py, "w", encoding="utf-8") as fh:
                    fh.write(
                        f'"""\nWeb Scraper Agent Starter — {clean_name}\n"""\n\n'
                        "import json\nimport os\n\n"
                        "def run_scraper():\n"
                        "    print('Starting web scraper for {clean_name}...')\n"
                        "    data = [{'item': 1, 'title': 'Sample Scraped Data', 'status': 'ready'}]\n"
                        "    with open('scraped_data.json', 'w', encoding='utf-8') as f:\n"
                        "        json.dump(data, f, indent=2)\n"
                        "    print('Saved results to scraped_data.json')\n\n"
                        "if __name__ == '__main__':\n"
                        "    run_scraper()\n"
                    )
                with open(os.path.join(proj_path, "requirements.txt"), "w", encoding="utf-8") as fh:
                    fh.write("requests>=2.31.0\nbeautifulsoup4>=4.12.0\n")

            elif runtime_lower in ("data_analyzer", "data", "analytics"):
                meta["runtime"] = "data_analyzer"
                starter_py = os.path.join(proj_path, "analyze_data.py")
                with open(starter_py, "w", encoding="utf-8") as fh:
                    fh.write(
                        f'"""\nData Analytics Agent Starter — {clean_name}\n"""\n\n'
                        "import json\nimport os\n\n"
                        "def analyze():\n"
                        "    print('Running analytics pipeline...')\n"
                        "    stats = {'metrics': {'rows': 100, 'accuracy': 0.98}, 'summary': 'Complete'}\n"
                        "    with open('report.json', 'w', encoding='utf-8') as f:\n"
                        "        json.dump(stats, f, indent=2)\n"
                        "    print('Generated report.json')\n\n"
                        "if __name__ == '__main__':\n"
                        "    analyze()\n"
                    )
                with open(os.path.join(proj_path, "requirements.txt"), "w", encoding="utf-8") as fh:
                    fh.write("pandas>=2.0.0\nmatplotlib>=3.7.0\n")

            elif runtime_lower in ("rest_api", "fastapi", "api"):
                meta["runtime"] = "rest_api"
                starter_py = os.path.join(proj_path, "server.py")
                with open(starter_py, "w", encoding="utf-8") as fh:
                    fh.write(
                        f'"""\nREST API Microservice Starter — {clean_name}\n"""\n\n'
                        "from http.server import HTTPServer, BaseHTTPRequestHandler\n"
                        "import json\n\n"
                        "class Handler(BaseHTTPRequestHandler):\n"
                        "    def do_GET(self):\n"
                        "        self.send_response(200)\n"
                        "        self.send_header('Content-Type', 'application/json')\n"
                        "        self.end_headers()\n"
                        "        self.wfile.write(json.dumps({'status': 'ok', 'project': '{clean_name}'}).encode('utf-8'))\n\n"
                        "def run():\n"
                        "    print('REST API ready on http://localhost:8080')\n\n"
                        "if __name__ == '__main__':\n"
                        "    run()\n"
                    )

            elif runtime_lower in ("html_dashboard", "dashboard", "web"):
                meta["runtime"] = "html_dashboard"
                with open(os.path.join(proj_path, "index.html"), "w", encoding="utf-8") as fh:
                    fh.write(
                        f"<!DOCTYPE html>\n<html lang=\"en\">\n<head>\n  <meta charset=\"UTF-8\">\n"
                        f"  <title>{clean_name} Dashboard</title>\n  <link rel=\"stylesheet\" href=\"style.css\">\n</head>\n"
                        f"<body>\n  <div class=\"container\">\n    <h1>{clean_name} Dashboard</h1>\n"
                        f"    <p>Powered by Text Surgeon Agent</p>\n    <div id=\"app\"></div>\n  </div>\n"
                        f"  <script src=\"app.js\"></script>\n</body>\n</html>\n"
                    )
                with open(os.path.join(proj_path, "style.css"), "w", encoding="utf-8") as fh:
                    fh.write("body { font-family: sans-serif; background: #0f172a; color: #f8fafc; margin: 0; padding: 2rem; }\n.container { max-width: 800px; margin: auto; }\n")
                with open(os.path.join(proj_path, "app.js"), "w", encoding="utf-8") as fh:
                    fh.write(f"console.log('{clean_name} initialized');\ndocument.getElementById('app').innerHTML = '<p>System status: Active</p>';\n")

            elif runtime_lower == "node":
                meta["runtime"] = "node"
                with open(os.path.join(proj_path, "index.js"), "w", encoding="utf-8") as fh:
                    fh.write(f'// Project: {clean_name}\nconsole.log("Hello from {clean_name}!");\n')
                with open(os.path.join(proj_path, "package.json"), "w", encoding="utf-8") as fh:
                    fh.write(json.dumps({"name": clean_name, "version": "1.0.0", "main": "index.js"}, indent=2) + "\n")

            else:
                meta["runtime"] = "python"
                starter_py = os.path.join(proj_path, "main.py")
                with open(starter_py, "w", encoding="utf-8") as fh:
                    fh.write(
                        f'"""\nProject: {clean_name}\nCreated: {now_iso}\n"""\n\n'
                        "def main():\n"
                        f'    print("Hello from {clean_name}!")\n\n'
                        'if __name__ == "__main__":\n'
                        "    main()\n"
                    )

            ws.save_project_meta(meta)

        return {
            "name": clean_name,
            "path": proj_path,
            "created_at": now_iso,
            "runtime": runtime,
            "description": description,
        }

    @classmethod
    def get_project_metadata(cls, path: str) -> Dict[str, Any]:
        """Returns metadata for a project directory."""
        ws = AgentWorkspace(path)
        return ws.get_project_meta()


class ProjectMemoryEngine:
    """Manages dual-tier Short-Term (session events) and Long-Term (synthesized knowledge) project memory."""

    def __init__(self, workspace_root: str) -> None:
        self.root = os.path.abspath(workspace_root)
        self.memory_dir = os.path.join(self.root, SURGEON_DIR_NAME, "memory")
        self.short_term_file = os.path.join(self.memory_dir, "short_term.json")
        self.long_term_file = os.path.join(self.memory_dir, "long_term.json")
        self.memory_md_file = os.path.join(self.memory_dir, "MEMORY.md")
        self.root_memory_md = os.path.join(self.root, SURGEON_DIR_NAME, "MEMORY.md")
        os.makedirs(self.memory_dir, exist_ok=True)

    def load_short_term(self) -> List[Dict[str, Any]]:
        """Loads short-term memory event list (recent conversation / task turns)."""
        if not os.path.isfile(self.short_term_file):
            return []
        try:
            with open(self.short_term_file, "r", encoding="utf-8") as fh:
                data = json.load(fh)
                return data if isinstance(data, list) else []
        except Exception:
            return []

    def record_short_term_event(
        self,
        event_type: str,
        goal: str = "",
        details: Optional[Dict[str, Any]] = None,
        max_events: int = 20,
    ) -> Dict[str, Any]:
        """Appends an event (goal, plan, execution result, user directive, error) to STM."""
        events = self.load_short_term()
        entry = {
            "id": f"evt_{int(time.time() * 1000)}",
            "timestamp": datetime.now().isoformat(),
            "type": event_type,
            "goal": goal,
            "details": details or {},
        }
        events.append(entry)
        events = events[-max_events:]
        try:
            with open(self.short_term_file, "w", encoding="utf-8") as fh:
                json.dump(events, fh, indent=2, ensure_ascii=False)
        except OSError:
            pass
        return entry

    def clear_short_term(self) -> None:
        """Clears short-term event buffer."""
        try:
            with open(self.short_term_file, "w", encoding="utf-8") as fh:
                json.dump([], fh, indent=2)
        except OSError:
            pass

    def load_long_term(self) -> Dict[str, Any]:
        """Loads persistent long-term knowledge items."""
        default_ltm: Dict[str, Any] = {
            "project_name": os.path.basename(self.root),
            "project_profile": {
                "tech_stack": [],
                "entrypoint": "",
                "env_vars_needed": [],
                "description": "",
            },
            "architecture_rules": [],
            "solved_issues_and_gotchas": [],
            "key_decisions": [],
            "user_preferences": [],
            "last_synthesized_at": None,
        }
        if not os.path.isfile(self.long_term_file):
            return default_ltm
        try:
            with open(self.long_term_file, "r", encoding="utf-8") as fh:
                data = json.load(fh)
                if isinstance(data, dict):
                    for k, v in default_ltm.items():
                        if k not in data:
                            data[k] = v
                    return data
                return default_ltm
        except Exception:
            return default_ltm

    def save_long_term(self, ltm: Dict[str, Any]) -> None:
        """Saves long-term memory and updates human-readable MEMORY.md."""
        ltm["last_modified"] = datetime.now().isoformat()
        try:
            with open(self.long_term_file, "w", encoding="utf-8") as fh:
                json.dump(ltm, fh, indent=2, ensure_ascii=False)
            self._render_memory_md(ltm)
        except OSError:
            pass

    def _render_memory_md(self, ltm: Dict[str, Any]) -> None:
        """Renders human-readable markdown for MEMORY.md."""
        lines = [
            f"# Project Memory & Knowledge Base — {ltm.get('project_name', os.path.basename(self.root))}",
            f"\n*Last updated: {ltm.get('last_modified', datetime.now().isoformat())}*",
            "",
            "## 1. Project Profile",
            f"- **Description**: {ltm.get('project_profile', {}).get('description') or 'N/A'}",
            f"- **Tech Stack**: {', '.join(ltm.get('project_profile', {}).get('tech_stack', [])) or 'Standard'}",
            f"- **Main Entrypoint**: `{ltm.get('project_profile', {}).get('entrypoint', 'main.py')}`",
        ]
        env_vars = ltm.get("project_profile", {}).get("env_vars_needed", [])
        if env_vars:
            lines.append(f"- **Required Environment Variables**: {', '.join(f'`{v}`' for v in env_vars)}")
        lines.append("")

        lines.append("## 2. Architecture & Conventions")
        arch = ltm.get("architecture_rules", [])
        if arch:
            for rule in arch:
                lines.append(f"- {rule}")
        else:
            lines.append("- *(No specific rules recorded yet)*")
        lines.append("")

        lines.append("## 3. Solved Issues & Gotchas")
        issues = ltm.get("solved_issues_and_gotchas", [])
        if issues:
            for item in issues:
                if isinstance(item, dict):
                    lines.append(f"- **{item.get('issue', 'Issue')}**: {item.get('solution', '')}")
                else:
                    lines.append(f"- {item}")
        else:
            lines.append("- *(No gotchas recorded yet)*")
        lines.append("")

        lines.append("## 4. Key Decisions & Preferences")
        decisions = ltm.get("key_decisions", []) + ltm.get("user_preferences", [])
        if decisions:
            for dec in decisions:
                lines.append(f"- {dec}")
        else:
            lines.append("- *(No decisions recorded yet)*")
        lines.append("")

        md_content = "\n".join(lines)
        for p in (self.memory_md_file, self.root_memory_md):
            try:
                with open(p, "w", encoding="utf-8") as fh:
                    fh.write(md_content)
            except OSError:
                pass

    def get_memory_summary_for_prompt(self, max_chars: int = 3500) -> str:
        """Produces an optimized context block for injection into LLM prompts."""
        ltm = self.load_long_term()
        stm = self.load_short_term()

        sections: List[str] = []

        # Long-term insights
        prof = ltm.get("project_profile", {})
        ltm_points: List[str] = []
        if prof.get("description"):
            ltm_points.append(f"Project Purpose: {prof['description']}")
        if prof.get("tech_stack"):
            ltm_points.append(f"Tech Stack: {', '.join(prof['tech_stack'])}")
        if prof.get("env_vars_needed"):
            ltm_points.append(f"Required .env Variables: {', '.join(prof['env_vars_needed'])}")

        arch = ltm.get("architecture_rules", [])
        if arch:
            ltm_points.append("Architecture Rules: " + " | ".join(arch[:5]))

        gotchas = ltm.get("solved_issues_and_gotchas", [])
        if gotchas:
            gotcha_strs = []
            for g in gotchas[:5]:
                if isinstance(g, dict):
                    gotcha_strs.append(f"{g.get('issue')}: {g.get('solution')}")
                else:
                    gotcha_strs.append(str(g))
            ltm_points.append("Past Solved Issues & Gotchas: " + " | ".join(gotcha_strs))

        if ltm_points:
            sections.append("### Persistent Knowledge (Long-Term Memory)\n" + "\n".join(f"- {p}" for p in ltm_points))

        # Short-term recent turns
        if stm:
            stm_lines = []
            for ev in stm[-4:]:
                g = ev.get("goal", "")
                typ = ev.get("type", "")
                det = ev.get("details", {})
                succ = det.get("success")
                status = " (Success)" if succ is True else (" (Failed)" if succ is False else "")
                stm_lines.append(f"- [{typ}{status}] {g[:140]}")
                if det.get("error"):
                    stm_lines.append(f"  Error encountered: {str(det['error'])[:120]}")
            if stm_lines:
                sections.append("### Recent Project Actions & Context (Short-Term Memory)\n" + "\n".join(stm_lines))

        if not sections:
            return ""

        full_block = "\n## PROJECT MEMORY & KNOWLEDGE BASE\n" + "\n\n".join(sections) + "\n"
        if len(full_block) > max_chars:
            full_block = full_block[:max_chars] + "\n...(memory truncated for brevity)\n"
        return full_block

    def synthesize_memory(
        self,
        llm_config: "LLMConfig",
        custom_instructions: str = "",
    ) -> Dict[str, Any]:
        """Synthesizes short-term memory + project state into long-term knowledge via LLM."""
        stm = self.load_short_term()
        ltm = self.load_long_term()
        scan = scan_workspace_tree(self.root, max_files=60)
        file_list = [f["path"] for f in scan.get("files", [])]

        prompt = (
            f"# PROJECT MEMORY SYNTHESIS PROTOCOL\n"
            f"You are the Text Surgeon Memory Synthesizer. Consolidate recent session events, project files, "
            f"and past lessons into persistent Long-Term Project Memory.\n\n"
            f"## PROJECT NAME: {os.path.basename(self.root)}\n"
            f"## EXISTING FILES IN WORKSPACE:\n{json.dumps(file_list, indent=2)}\n\n"
            f"## CURRENT LONG-TERM MEMORY:\n{json.dumps(ltm, indent=2)}\n\n"
            f"## RECENT SHORT-TERM EVENTS & GOALS:\n{json.dumps(stm, indent=2)}\n\n"
        )
        if custom_instructions:
            prompt += f"## SPECIAL DIRECTIVES:\n{custom_instructions}\n\n"

        prompt += (
            "## OUTPUT FORMAT\n"
            "Return ONLY a valid JSON object matching this exact schema (no markdown fences, no extra text):\n"
            "{\n"
            '  "project_profile": {\n'
            '    "description": "Concise summary of what this project does",\n'
            '    "tech_stack": ["python", "python-pptx", "python-telegram-bot"],\n'
            '    "entrypoint": "main.py or bot.py",\n'
            '    "env_vars_needed": ["TELEGRAM_BOT_TOKEN"]\n'
            "  },\n"
            '  "architecture_rules": [\n'
            '    "Rule 1: Use python-dotenv to load environment variables",\n'
            '    "Rule 2: Keep presentations in root and bot code modular"\n'
            "  ],\n"
            '  "solved_issues_and_gotchas": [\n'
            '    {\n'
            '      "issue": "Telegram Bot InvalidToken error",\n'
            '      "solution": "Requires valid TELEGRAM_BOT_TOKEN in .env; placeholder will fail auth."\n'
            "    }\n"
            "  ],\n"
            '  "key_decisions": [\n'
            '    "Decision 1: User requested 10 slides with Persian text on last slide"\n'
            "  ],\n"
            '  "user_preferences": [\n'
            '    "Preference 1: Use clean async handlers and proper logging"\n'
            "  ]\n"
            "}\n"
        )

        raw_reply = LLMClient.send_prompt(
            prompt=prompt,
            config=llm_config,
            system_prompt="You are a precise JSON memory extraction and consolidation engine.",
        )

        cleaned = raw_reply.strip()
        if "```json" in cleaned:
            cleaned = cleaned.split("```json", 1)[1].split("```", 1)[0].strip()
        elif "```" in cleaned:
            cleaned = cleaned.split("```", 1)[1].split("```", 1)[0].strip()

        try:
            parsed = json.loads(cleaned)
        except Exception:
            m = re.search(r"\{.*\}", cleaned, re.DOTALL)
            if m:
                parsed = json.loads(m.group(0))
            else:
                raise AgentError(f"Failed to parse LLM memory synthesis response as JSON: {raw_reply[:300]}")

        ltm["project_profile"] = parsed.get("project_profile", ltm.get("project_profile", {}))
        ltm["architecture_rules"] = parsed.get("architecture_rules", ltm.get("architecture_rules", []))
        ltm["solved_issues_and_gotchas"] = parsed.get("solved_issues_and_gotchas", ltm.get("solved_issues_and_gotchas", []))
        ltm["key_decisions"] = parsed.get("key_decisions", ltm.get("key_decisions", []))
        ltm["user_preferences"] = parsed.get("user_preferences", ltm.get("user_preferences", []))
        ltm["last_synthesized_at"] = datetime.now().isoformat()

        self.save_long_term(ltm)
        return ltm


class AgentWorkspace:
    """Manages files, structured .surgeon backups, logs, and execution within a project."""

    def __init__(self, workspace_root: str) -> None:
        self.root = os.path.abspath(workspace_root)
        os.makedirs(self.root, exist_ok=True)

        self.surgeon_dir = os.path.join(self.root, SURGEON_DIR_NAME)
        self.backups_dir = os.path.join(self.surgeon_dir, "backups")
        self.logs_dir = os.path.join(self.surgeon_dir, "logs")
        self.meta_file = os.path.join(self.surgeon_dir, "project.json")
        self.memory = ProjectMemoryEngine(self.root)

        os.makedirs(self.backups_dir, exist_ok=True)
        os.makedirs(self.logs_dir, exist_ok=True)

    def get_project_meta(self) -> Dict[str, Any]:
        """Reads project metadata from .surgeon/project.json."""
        if not os.path.isfile(self.meta_file):
            return {
                "name": os.path.basename(self.root),
                "created_at": datetime.now().isoformat(),
                "last_modified": datetime.now().isoformat(),
                "runtime": "python",
                "goals": [],
            }
        try:
            with open(self.meta_file, "r", encoding="utf-8") as fh:
                return json.load(fh)
        except Exception:
            return {"name": os.path.basename(self.root)}

    def save_project_meta(self, meta: Dict[str, Any]) -> None:
        """Saves project metadata to .surgeon/project.json."""
        meta["last_modified"] = datetime.now().isoformat()
        try:
            with open(self.meta_file, "w", encoding="utf-8") as fh:
                json.dump(meta, fh, indent=2, ensure_ascii=False)
        except OSError:
            pass

    def scan(self) -> Dict[str, Any]:
        """Scans the workspace directory."""
        return scan_workspace_tree(self.root)

    def preview_plan(self, plan: AgentPlan) -> Dict[str, Any]:
        """Previews changes without modifying disk.

        Returns:
            Dict of created, modified (with diffs), deleted files, and commands to run.
        """
        changes: List[Dict[str, Any]] = []

        for action in plan.file_actions:
            safe_target = resolve_safe_workspace_path(self.root, action.path)
            rel_display = os.path.relpath(safe_target, self.root).replace("\\", "/")
            exists = os.path.isfile(safe_target)

            if action.action_type == "delete":
                changes.append(
                    {
                        "path": rel_display,
                        "action": "delete",
                        "exists": exists,
                        "diff": f"- (delete {rel_display})" if exists else "(already missing)",
                    }
                )
            elif action.action_type in ("create", "modify", "overwrite"):
                old_text = ""
                if exists:
                    try:
                        with open(safe_target, "r", encoding="utf-8", errors="replace") as fh:
                            old_text = fh.read()
                    except OSError:
                        old_text = ""
                diff = "".join(
                    difflib.unified_diff(
                        old_text.splitlines(keepends=True),
                        action.content.splitlines(keepends=True),
                        fromfile=f"a/{rel_display}" if exists else "/dev/null",
                        tofile=f"b/{rel_display}",
                    )
                )
                changes.append(
                    {
                        "path": rel_display,
                        "action": "overwrite" if exists else "create",
                        "exists": exists,
                        "lines": action.content.count("\n") + (1 if action.content else 0),
                        "chars": len(action.content),
                        "diff": diff or "(identical content)",
                        "preview": action.content[:600] + ("..." if len(action.content) > 600 else ""),
                    }
                )
            elif action.action_type == "edit":
                if not exists:
                    changes.append(
                        {
                            "path": rel_display,
                            "action": "error",
                            "error": f"File '{rel_display}' does not exist to edit.",
                        }
                    )
                    continue
                try:
                    with open(safe_target, "r", encoding="utf-8") as fh:
                        old_text = fh.read()
                    new_text, applied, skipped = apply_edit_ops(old_text, action.edit_ops)
                    diff = "".join(
                        difflib.unified_diff(
                            old_text.splitlines(keepends=True),
                            new_text.splitlines(keepends=True),
                            fromfile=f"a/{rel_display}",
                            tofile=f"b/{rel_display}",
                        )
                    )
                    changes.append(
                        {
                            "path": rel_display,
                            "action": "edit",
                            "exists": True,
                            "splices": len(applied),
                            "diff": diff,
                        }
                    )
                except SelectionError as exc:
                    changes.append(
                        {
                            "path": rel_display,
                            "action": "edit_refusal",
                            "error": str(exc),
                            "details": getattr(exc, "details", {}),
                        }
                    )

        return {
            "root": self.root,
            "explanation": plan.explanation,
            "changes": changes,
            "setup_commands": plan.setup_commands,
            "run_command": plan.run_command,
            "expected_artifacts": plan.expected_artifacts,
        }

    def list_backups(self) -> List[Dict[str, Any]]:
        """Lists all organized snapshot backups in .surgeon/backups/."""
        backups: List[Dict[str, Any]] = []
        if not os.path.isdir(self.backups_dir):
            return backups

        try:
            entries = sorted(os.listdir(self.backups_dir), reverse=True)
        except OSError:
            return backups

        for entry in entries:
            snap_dir = os.path.join(self.backups_dir, entry)
            if not os.path.isdir(snap_dir):
                continue
            meta_path = os.path.join(snap_dir, "meta.json")
            if os.path.isfile(meta_path):
                try:
                    with open(meta_path, "r", encoding="utf-8") as fh:
                        info = json.load(fh)
                        info["backup_id"] = entry
                        backups.append(info)
                        continue
                except Exception:
                    pass

            # Fallback if meta.json missing
            files_in_snap = []
            for r, _, f_list in os.walk(snap_dir):
                for f in f_list:
                    if f != "meta.json":
                        files_in_snap.append(os.path.relpath(os.path.join(r, f), snap_dir).replace("\\", "/"))

            backups.append(
                {
                    "backup_id": entry,
                    "timestamp": entry,
                    "explanation": "Snapshot backup",
                    "files": files_in_snap,
                }
            )

        return backups

    def restore_backup(self, backup_id: str, target_file: Optional[str] = None) -> Dict[str, Any]:
        """Restores all files or a specific file from a snapshot backup.

        Args:
            backup_id: The snapshot ID (subfolder in .surgeon/backups/).
            target_file: Optional relative file path to restore. If None, restores all.

        Returns:
            Dict with success status and list of restored files.
        """
        clean_id = os.path.basename(backup_id)
        snap_dir = os.path.join(self.backups_dir, clean_id)
        if not os.path.isdir(snap_dir):
            raise AgentError(f"Backup snapshot '{backup_id}' does not exist.")

        restored_files: List[str] = []

        if target_file:
            safe_target_rel = target_file.replace("\\", "/").strip("/")
            snap_src = os.path.join(snap_dir, safe_target_rel)
            if not os.path.isfile(snap_src):
                raise AgentError(f"File '{target_file}' not found in backup '{backup_id}'.")
            dest_target = resolve_safe_workspace_path(self.root, safe_target_rel)
            os.makedirs(os.path.dirname(dest_target), exist_ok=True)
            shutil.copy2(snap_src, dest_target)
            restored_files.append(safe_target_rel)
        else:
            for root_dir, _, files in os.walk(snap_dir):
                for f in files:
                    if f == "meta.json":
                        continue
                    src_file = os.path.join(root_dir, f)
                    rel_p = os.path.relpath(src_file, snap_dir).replace("\\", "/")
                    dest_file = resolve_safe_workspace_path(self.root, rel_p)
                    os.makedirs(os.path.dirname(dest_file), exist_ok=True)
                    shutil.copy2(src_file, dest_file)
                    restored_files.append(rel_p)

        return {
            "success": True,
            "backup_id": clean_id,
            "restored_files": restored_files,
        }

    def apply_plan(self, plan: AgentPlan, backup: bool = True) -> Dict[str, Any]:
        """Applies all file actions transactionally inside workspace with organized backups.

        Returns:
            Summary of written/edited/deleted files, backup ID, and root path.
        """
        written_files: List[str] = []
        backup_id = f"{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}_{uuid.uuid4().hex[:6]}"
        snap_dir = os.path.join(self.backups_dir, backup_id)
        backed_up_files: List[str] = []

        try:
            # 1. Pre-backup affected files into .surgeon/backups/<backup_id>/
            if backup:
                for action in plan.file_actions:
                    safe_target = resolve_safe_workspace_path(self.root, action.path)
                    if os.path.isfile(safe_target):
                        rel_path = os.path.relpath(safe_target, self.root).replace("\\", "/")
                        snap_target = os.path.join(snap_dir, rel_path)
                        os.makedirs(os.path.dirname(snap_target), exist_ok=True)
                        shutil.copy2(safe_target, snap_target)
                        backed_up_files.append(rel_path)

                if backed_up_files or plan.file_actions:
                    os.makedirs(snap_dir, exist_ok=True)
                    meta = {
                        "backup_id": backup_id,
                        "timestamp": datetime.now().isoformat(),
                        "explanation": plan.explanation or "AI edit plan",
                        "files": backed_up_files,
                        "actions": [{"path": a.path, "action": a.action_type} for a in plan.file_actions],
                    }
                    with open(os.path.join(snap_dir, "meta.json"), "w", encoding="utf-8") as fh:
                        json.dump(meta, fh, indent=2, ensure_ascii=False)

            # 2. Apply modifications atomically
            for action in plan.file_actions:
                safe_target = resolve_safe_workspace_path(self.root, action.path)
                parent_dir = os.path.dirname(safe_target)
                os.makedirs(parent_dir, exist_ok=True)

                if action.action_type == "delete":
                    if os.path.isfile(safe_target):
                        os.unlink(safe_target)
                        written_files.append(action.path)

                elif action.action_type in ("create", "modify", "overwrite"):
                    tmp_target = safe_target + ".tmp"
                    with open(tmp_target, "w", encoding="utf-8", newline="\n") as fh:
                        fh.write(action.content)
                        fh.flush()
                        os.fsync(fh.fileno())
                    os.replace(tmp_target, safe_target)
                    written_files.append(action.path)

                elif action.action_type == "edit":
                    if not os.path.isfile(safe_target):
                        raise AgentError(f"Cannot edit non-existent file '{action.path}'.")
                    with open(safe_target, "r", encoding="utf-8") as fh:
                        old_text = fh.read()
                    new_text, applied, skipped = apply_edit_ops(old_text, action.edit_ops)
                    tmp_target = safe_target + ".tmp"
                    with open(tmp_target, "w", encoding="utf-8", newline="\n") as fh:
                        fh.write(new_text)
                        fh.flush()
                        os.fsync(fh.fileno())
                    os.replace(tmp_target, safe_target)
                    written_files.append(action.path)

            # Update project metadata
            proj_meta = self.get_project_meta()
            if plan.explanation:
                goals = proj_meta.get("goals", [])
                goals.append({"timestamp": datetime.now().isoformat(), "goal": plan.explanation})
                proj_meta["goals"] = goals[-20:]  # keep recent 20
            self.save_project_meta(proj_meta)

        except Exception as exc:
            # Revert from snapshot on error
            if backed_up_files and os.path.isdir(snap_dir):
                for rel_f in backed_up_files:
                    src_bak = os.path.join(snap_dir, rel_f)
                    dest_f = resolve_safe_workspace_path(self.root, rel_f)
                    try:
                        if os.path.isfile(src_bak):
                            shutil.copy2(src_bak, dest_f)
                    except OSError:
                        pass
            raise AgentError(f"Failed to apply plan to workspace: {exc}") from exc

        return {
            "success": True,
            "written_files": written_files,
            "backup_id": backup_id if backup else None,
            "root": self.root,
        }

    def install_dependencies(
        self,
        packages: Sequence[str],
        timeout: int = 300,
    ) -> ExecutionResult:
        """Installs Python packages via pip safely inside workspace.

        Args:
            packages: List/tuple of package names or requirement strings.
            timeout: Max seconds to wait for pip installation.

        Returns:
            ExecutionResult containing pip output and exit status.
        """
        clean_pkgs = [p.strip() for p in packages if p and p.strip()]
        if not clean_pkgs:
            return ExecutionResult(
                command="",
                exit_code=0,
                stdout="No packages to install.",
                stderr="",
                duration_sec=0.0,
                success=True,
            )
        py_bin = sys.executable or "python"
        cmd = f'"{py_bin}" -m pip install {" ".join(clean_pkgs)}'
        return self.execute_command(cmd, timeout=timeout)

    def execute_command(
        self,
        command: str,
        timeout: int = DEFAULT_EXECUTION_TIMEOUT,
        expected_artifacts: Optional[List[str]] = None,
        timeout_sec: Optional[int] = None,
    ) -> ExecutionResult:
        """Executes a command inside the project directory and tracks results."""
        if timeout_sec is not None:
            timeout = timeout_sec
        clean_cmd = command.strip()
        if not clean_cmd:
            return ExecutionResult(
                command="",
                exit_code=0,
                stdout="No command executed.",
                stderr="",
                duration_sec=0.0,
                success=True,
            )

        before_files = set()
        for root, _, files in os.walk(self.root):
            # Skip .surgeon directory from artifact detection
            if os.path.relpath(root, self.root).startswith(SURGEON_DIR_NAME):
                continue
            for f in files:
                before_files.add(os.path.relpath(os.path.join(root, f), self.root))

        start_time = time.perf_counter()
        timed_out = False

        env = os.environ.copy()
        env["PYTHONUNBUFFERED"] = "1"
        existing_py = env.get("PYTHONPATH", "")
        env["PYTHONPATH"] = f"{self.root}{os.pathsep}{existing_py}" if existing_py else self.root
        # Inject project .env variables
        project_env = EnvManager.load(self.root)
        env.update(project_env)

        try:
            proc = subprocess.run(
                clean_cmd,
                cwd=self.root,
                shell=True,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=timeout,
                env=env,
            )
            duration = time.perf_counter() - start_time
            exit_code = proc.returncode
            stdout_str = proc.stdout or ""
            stderr_str = proc.stderr or ""

        except subprocess.TimeoutExpired as exc:
            duration = time.perf_counter() - start_time
            timed_out = True
            exit_code = -1
            stdout_str = (exc.stdout or "").decode("utf-8", "replace") if isinstance(exc.stdout, bytes) else (exc.stdout or "")
            stderr_str = f"Execution timed out after {timeout} seconds."

        except Exception as exc:
            duration = time.perf_counter() - start_time
            exit_code = 1
            stdout_str = ""
            stderr_str = f"Failed to launch command '{clean_cmd}': {exc}"

        after_files = set()
        for root, _, files in os.walk(self.root):
            if os.path.relpath(root, self.root).startswith(SURGEON_DIR_NAME):
                continue
            for f in files:
                after_files.add(os.path.relpath(os.path.join(root, f), self.root))

        new_artifacts = sorted(list(after_files - before_files))

        if expected_artifacts:
            for art in expected_artifacts:
                art_clean = art.strip().replace("\\", "/")
                safe_art = resolve_safe_workspace_path(self.root, art_clean)
                if os.path.exists(safe_art) and art_clean not in new_artifacts:
                    new_artifacts.append(art_clean)

        res = ExecutionResult(
            command=clean_cmd,
            exit_code=exit_code,
            stdout=stdout_str,
            stderr=stderr_str,
            duration_sec=duration,
            artifacts_found=new_artifacts,
            timed_out=timed_out,
            success=(exit_code == 0),
        )

        # Save structured log to .surgeon/logs/
        try:
            log_id = datetime.now().strftime("%Y%m%d_%H%M%S")
            log_file = os.path.join(self.logs_dir, f"run_{log_id}.json")
            with open(log_file, "w", encoding="utf-8") as fh:
                json.dump(
                    {
                        "command": clean_cmd,
                        "exit_code": exit_code,
                        "duration_sec": duration,
                        "stdout": stdout_str,
                        "stderr": stderr_str,
                        "artifacts": new_artifacts,
                        "timestamp": datetime.now().isoformat(),
                        "success": exit_code == 0,
                    },
                    fh,
                    indent=2,
                    ensure_ascii=False,
                )
        except OSError:
            pass

        return res

    def read_file(self, relative_path: str, max_bytes: int = 10 * 1024 * 1024) -> str:
        """Reads file content safely from the workspace."""
        safe_path = resolve_safe_workspace_path(self.root, relative_path)
        if not os.path.isfile(safe_path):
            raise AgentError(f"File '{relative_path}' does not exist.")
        if os.path.getsize(safe_path) > max_bytes:
            raise AgentError(f"File '{relative_path}' exceeds maximum viewable size ({max_bytes} bytes).")
        with open(safe_path, "r", encoding="utf-8", errors="replace") as fh:
            return fh.read()

    def write_file(self, relative_path: str, content: str) -> Dict[str, Any]:
        """Writes or updates a file inside workspace atomically with snapshot backup."""
        safe_path = resolve_safe_workspace_path(self.root, relative_path)
        os.makedirs(os.path.dirname(safe_path), exist_ok=True)
        tmp_target = safe_path + ".tmp"
        with open(tmp_target, "w", encoding="utf-8", newline="\n") as fh:
            fh.write(content)
            fh.flush()
            os.fsync(fh.fileno())
        os.replace(tmp_target, safe_path)
        return {
            "success": True,
            "path": relative_path,
            "size": os.path.getsize(safe_path),
        }

    def delete_file(self, relative_path: str) -> bool:
        """Deletes a file safely from workspace."""
        safe_path = resolve_safe_workspace_path(self.root, relative_path)
        if os.path.isfile(safe_path):
            os.unlink(safe_path)
            return True
        return False

    def get_env(self) -> Dict[str, str]:
        """Loads workspace .env variables."""
        return EnvManager.load(self.root)

    def set_env(self, env_vars: Dict[str, str]) -> None:
        """Saves workspace .env variables."""
        EnvManager.save(self.root, env_vars)

    def export_zip(self) -> bytes:
        """Exports the entire workspace as a zip archive."""
        return export_project_zip(self.root)

    def record_round(self, round_data: Dict[str, Any]) -> None:
        """Records an agent round (prompt, plan, execution logs) into .surgeon/rounds.json."""
        rounds_file = os.path.join(self.surgeon_dir, "rounds.json")
        rounds: List[Dict[str, Any]] = []
        if os.path.isfile(rounds_file):
            try:
                with open(rounds_file, "r", encoding="utf-8") as fh:
                    rounds = json.load(fh)
            except Exception:
                rounds = []
        rounds.append(round_data)
        # Keep recent 50 rounds
        rounds = rounds[-50:]
        try:
            with open(rounds_file, "w", encoding="utf-8") as fh:
                json.dump(rounds, fh, indent=2, ensure_ascii=False)
        except OSError:
            pass

    def get_round_history(self) -> List[Dict[str, Any]]:
        """Returns the history of agent rounds recorded in this workspace."""
        rounds_file = os.path.join(self.surgeon_dir, "rounds.json")
        if not os.path.isfile(rounds_file):
            return []
        try:
            with open(rounds_file, "r", encoding="utf-8") as fh:
                return json.load(fh)
        except Exception:
            return []

# --------------------------------------------------------------------------- #
# Extensible Skills Subsystem
# --------------------------------------------------------------------------- #


@dataclass
class Skill:
    """A modular domain-specific capability for Text Surgeon Agent."""

    name: str
    title: str
    description: str
    keywords: List[str] = field(default_factory=list)
    packages: List[str] = field(default_factory=list)
    content: str = ""
    file_path: Optional[str] = None
    source: str = "builtin"  # 'builtin' | 'global' | 'project' | 'custom'

    def to_dict(self) -> Dict[str, Any]:
        """Serializes skill metadata for JSON APIs."""
        return {
            "name": self.name,
            "title": self.title,
            "description": self.description,
            "keywords": self.keywords,
            "packages": self.packages,
            "content": self.content,
            "file_path": self.file_path,
            "source": self.source,
        }


class SkillManager:
    """Discovers, parses, auto-detects, and imports modular Agent skills."""

    APP_DIR: Final[str] = os.path.dirname(os.path.abspath(__file__))
    GLOBAL_SKILLS_DIR: Final[str] = os.path.expanduser("~/.gemini/config/skills")

    @classmethod
    def parse_skill_markdown(
        cls,
        text: str,
        file_path: Optional[str] = None,
        source: str = "custom",
    ) -> Skill:
        """Parses a SKILL.md file containing YAML frontmatter and markdown body."""
        clean = text.strip()
        name = "custom_skill"
        title = "Custom Skill"
        description = ""
        keywords: List[str] = []
        packages: List[str] = []
        body = clean

        if clean.startswith("---"):
            parts = clean.split("---", 2)
            if len(parts) >= 3:
                frontmatter_text = parts[1].strip()
                body = parts[2].strip()
                for line in frontmatter_text.splitlines():
                    line = line.strip()
                    if not line or line.startswith("#"):
                        continue
                    if ":" in line:
                        k, v = line.split(":", 1)
                        k = k.strip().lower()
                        v = v.strip()
                        if k == "name":
                            name = v.strip("\"'")
                        elif k == "title":
                            title = v.strip("\"'")
                        elif k == "description":
                            description = v.strip("\"'")
                        elif k == "keywords":
                            # Parse [kw1, kw2] or comma list
                            raw_kws = v.strip("[] \t\r\n")
                            keywords = [kw.strip("\"' ") for kw in raw_kws.split(",") if kw.strip("\"' ")]
                        elif k == "packages":
                            raw_pkgs = v.strip("[] \t\r\n")
                            packages = [pkg.strip("\"' ") for pkg in raw_pkgs.split(",") if pkg.strip("\"' ")]

        if file_path and (not name or name == "custom_skill"):
            parent_dir = os.path.basename(os.path.dirname(file_path))
            if parent_dir and parent_dir != "skills":
                name = parent_dir
            else:
                base = os.path.splitext(os.path.basename(file_path))[0]
                name = base if base != "SKILL" else "skill"

        if not title or title == "Custom Skill":
            title = name.replace("_", " ").title()

        return Skill(
            name=name,
            title=title,
            description=description or f"Skill for {title}",
            keywords=keywords,
            packages=packages,
            content=body,
            file_path=file_path,
            source=source,
        )

    @classmethod
    def get_skills_dirs(cls, workspace_root: Optional[str] = None) -> List[Tuple[str, str]]:
        """Returns the search directories for skills (path, source_tag)."""
        dirs = [
            (os.path.join(cls.APP_DIR, "skills"), "builtin"),
            (cls.GLOBAL_SKILLS_DIR, "global"),
        ]
        if workspace_root:
            dirs.append((os.path.join(workspace_root, "skills"), "project"))
        return dirs

    @classmethod
    def discover_skills(cls, workspace_root: Optional[str] = None) -> List[Skill]:
        """Discovers all available skills across builtin, global, and project directories."""
        discovered: Dict[str, Skill] = {}
        for search_dir, source in cls.get_skills_dirs(workspace_root):
            if not os.path.isdir(search_dir):
                continue
            for item in os.listdir(search_dir):
                item_path = os.path.join(search_dir, item)
                skill_file = None
                if os.path.isdir(item_path):
                    cand = os.path.join(item_path, "SKILL.md")
                    if os.path.isfile(cand):
                        skill_file = cand
                elif item.endswith(".md"):
                    skill_file = item_path

                if skill_file and os.path.isfile(skill_file):
                    try:
                        with open(skill_file, "r", encoding="utf-8", errors="replace") as fh:
                            content = fh.read()
                        skill = cls.parse_skill_markdown(content, file_path=skill_file, source=source)
                        # Project/custom skills override builtin skills with same name
                        discovered[skill.name] = skill
                    except OSError:
                        continue
        return list(discovered.values())

    @classmethod
    def get_skill(cls, name: str, workspace_root: Optional[str] = None) -> Optional[Skill]:
        """Fetches a skill by name."""
        clean_name = name.lower().strip()
        for skill in cls.discover_skills(workspace_root):
            if skill.name.lower() == clean_name:
                return skill
        return None

    @classmethod
    def detect_skills(cls, goal: str, workspace_root: Optional[str] = None) -> List[Skill]:
        """Matches keywords in user prompt/goal to discover relevant skills."""
        clean_goal = goal.lower()
        matched: List[Skill] = []
        for skill in cls.discover_skills(workspace_root):
            # Check skill name in goal
            if skill.name.lower() in clean_goal:
                matched.append(skill)
                continue
            # Check keywords
            for kw in skill.keywords:
                if kw.lower() in clean_goal:
                    matched.append(skill)
                    break
        return matched

    @classmethod
    def import_from_url(cls, url: str, target_dir: Optional[str] = None) -> Skill:
        """Downloads a SKILL.md file from a URL and installs it into skills directory."""
        clean_url = url.strip()
        req = urllib.request.Request(
            clean_url,
            headers={"User-Agent": "TextSurgeon-SkillImporter/2.0"},
        )
        with urllib.request.urlopen(req, timeout=15) as resp:
            raw_data = resp.read().decode("utf-8", errors="replace")

        skill = cls.parse_skill_markdown(raw_data, source="custom")
        save_dir = target_dir or os.path.join(cls.APP_DIR, "skills")
        dest_folder = os.path.join(save_dir, skill.name)
        os.makedirs(dest_folder, exist_ok=True)
        dest_file = os.path.join(dest_folder, "SKILL.md")
        with open(dest_file, "w", encoding="utf-8", newline="\n") as fh:
            fh.write(raw_data)
        skill.file_path = dest_file
        return skill

    @classmethod
    def create_skill(
        cls,
        name: str,
        title: str,
        description: str,
        content: str,
        keywords: Sequence[str] = (),
        packages: Sequence[str] = (),
        target_dir: Optional[str] = None,
    ) -> Skill:
        """Creates and saves a new skill."""
        clean_name = re.sub(r"[^a-zA-Z0-9_\-]", "_", name.strip().lower())
        if not clean_name:
            raise ValueError("Invalid skill name.")
        save_dir = target_dir or os.path.join(cls.APP_DIR, "skills")
        dest_folder = os.path.join(save_dir, clean_name)
        os.makedirs(dest_folder, exist_ok=True)
        dest_file = os.path.join(dest_folder, "SKILL.md")

        kws_str = ", ".join(keywords)
        pkgs_str = ", ".join(packages)
        md_text = f"""---
name: {clean_name}
title: {title or clean_name.replace('_', ' ').title()}
description: {description}
keywords: [{kws_str}]
packages: [{pkgs_str}]
---

{content.strip()}
"""
        with open(dest_file, "w", encoding="utf-8", newline="\n") as fh:
            fh.write(md_text)

        return Skill(
            name=clean_name,
            title=title or clean_name.replace("_", " ").title(),
            description=description,
            keywords=list(keywords),
            packages=list(packages),
            content=content.strip(),
            file_path=dest_file,
            source="custom",
        )

    @classmethod
    def create_custom_skill(
        cls,
        workspace_root: str,
        name: str,
        title: str = "",
        description: str = "",
        keywords: Sequence[str] = (),
        packages: Sequence[str] = (),
        content: str = "",
    ) -> Skill:
        """Creates a custom skill stored inside workspace/skills."""
        target_dir = os.path.join(workspace_root, "skills")
        return cls.create_skill(
            name=name,
            title=title,
            description=description,
            content=content,
            keywords=keywords,
            packages=packages,
            target_dir=target_dir,
        )

    @classmethod
    def match_skills_to_goal(cls, goal: str, workspace_root: Optional[str] = None) -> List[Skill]:
        """Alias for detect_skills."""
        return cls.detect_skills(goal, workspace_root)



# --------------------------------------------------------------------------- #
# Agent Prompt Builders
# --------------------------------------------------------------------------- #


_AGENT_PROTOCOL_RULES: Final[str] = """\
## AGENT PROTOCOL v1.0 RULES

You are operating as an autonomous software engineer and project agent.
Your mission is to write complete, robust, ready-to-run code to achieve the user's goal.

### ADAPTIVE TASK SIZING & NON-REDUNDANT MULTI-STEP PROTOCOL
1. **Single-Step Tasks** (e.g., creating a PowerPoint, writing a script, creating a microservice, single-file edits):
   - Deliver the complete, production-ready solution in a SINGLE comprehensive round.
   - Provide all necessary @@FILE blocks and the execution @@RUN command immediately.
   - Do NOT split simple tasks into artificial rounds.

2. **Investigative & Complex Multi-Step Tasks** (e.g., "diagnose my PC", "analyze system bottlenecks", "multi-phase refactoring", "debug why network requests fail"):
   - Use an adaptive multi-phase approach:
     * Phase 1: Write and run a diagnostic probe / diagnostic script to gather system info or error metrics.
     * Phase 2: In the subsequent round, analyze the diagnostic logs, formulate the fix, and apply the final solution.
   - You may emit `@@STEP <current_step>/<total_steps>: <step_description>` in your explanation.

3. **Zero Redundancy**: When the objective is achieved with exit code 0 and required outputs produced, stop immediately.

### MANDATORY RESPONSE FORMAT

Your response must contain:
1. An <EXPLANATION> section explaining your architecture, libraries chosen, and design.
2. @@COMMAND lines for any packages/dependencies that need installation (e.g. `pip install python-pptx` or `npm install express`).
3. @@FILE blocks for creating new files or full rewrites, OR @@EDIT blocks for precision surgical in-place edits on existing files.
4. @@RUN line indicating the exact command to execute to run the project or test.
5. Optional @@ARTIFACT lines for expected deliverables (e.g. `presentation.pptx`, `index.html`).

#### Block Grammar:

<EXPLANATION>
Brief overview of your solution, dependencies, and file layout.
</EXPLANATION>

@@COMMAND pip install <package_name>

@@FILE path/to/new_file.py
<<<
# Full, complete code for new files or total rewrites.
>>>

@@EDIT path/to/existing_file.py anchor
START-ANCHOR: first 5-10 words of section to change
END-ANCHOR: last 5-10 words of section to change
<<<
# Surgical in-place replacement for just this block.
>>>

@@RUN python path/to/main.py
@@ARTIFACT output_file.pptx

### HARD CODING RULES:
1. NO TRUNCATION / NO PLACEHOLDERS: Never write "# ... rest of code goes here ...". Write complete, runnable code.
2. PURE & SAFE PATHS: Use relative file paths inside the workspace.
3. ERROR HANDLING: Ensure scripts handle edge cases, create required folders if needed, and print informative progress messages to stdout. If an error is fatal, do NOT swallow it with exit code 0; raise or exit(1) so the repair loop can heal it.
4. ENVIRONMENT CONFIG: Always load environment variables (e.g. `python-dotenv`) for credentials/tokens instead of hardcoding placeholders.
5. DELIMITERS: Write `<<<` and `>>>` on their own lines as plain ASCII.
"""


class AgentPromptBuilder:
    """Generates Agent Protocol prompts for new tasks, follow-ups, and verification loops."""

    @staticmethod
    def build_task_prompt(
        goal: str,
        workspace_root: str,
        runtime: str = "python",
        include_workspace_files: bool = True,
        include_memory: bool = True,
        skill_names: Optional[Sequence[str]] = None,
        auto_detect_skills: bool = True,
    ) -> str:
        """Builds the initial Agent Protocol prompt to send to the LLM.

        Args:
            goal: The user's task description (e.g. 'Create a 10-slide PowerPoint...').
            workspace_root: Target project folder.
            runtime: Preferred language environment ('python', 'node', 'shell', etc.).
            include_workspace_files: Whether to inspect and embed existing workspace files.
            include_memory: Whether to attach project Short-Term and Long-Term Memory.
            skill_names: Explicit skill names to apply.
            auto_detect_skills: If True and skill_names is empty, matches keywords automatically.

        Returns:
            The complete Agent Protocol prompt.
        """
        scan = scan_workspace_tree(workspace_root) if include_workspace_files else {"files": []}
        files = scan.get("files", [])

        # Resolve active skills
        active_skills: List[Skill] = []
        if skill_names:
            for sname in skill_names:
                sk = SkillManager.get_skill(sname, workspace_root)
                if sk and sk not in active_skills:
                    active_skills.append(sk)
        elif auto_detect_skills:
            active_skills = SkillManager.detect_skills(goal, workspace_root)

        blocks: List[str] = [
            f"# SURGEON AGENT PROTOCOL v{AGENT_PROTOCOL_VERSION} — PROJECT BUILDER\n",
            _AGENT_PROTOCOL_RULES,
            "\n## TARGET ENVIRONMENT & RUNTIME\n",
            f"Primary Language/Runtime: {runtime.upper()}\n",
            f"Project Directory: {os.path.basename(workspace_root)}\n",
            "\n## USER GOAL / TASK REQUEST\n\n",
            goal.strip(),
            "\n",
        ]

        if include_memory:
            mem_engine = ProjectMemoryEngine(workspace_root)
            mem_block = mem_engine.get_memory_summary_for_prompt()
            if mem_block:
                blocks.append(mem_block)

        if active_skills:
            blocks.append("\n## SPECIALIZED SKILLS & DOMAIN GUIDELINES\n")
            blocks.append("Apply the following specialized skills and guidelines to fulfill this task:\n")
            for sk in active_skills:
                blocks.append(f"\n### SKILL: {sk.title} (`{sk.name}`)\n")
                if sk.packages:
                    blocks.append(f"**Required Packages**: `{', '.join(sk.packages)}`\n")
                    for pkg in sk.packages:
                        blocks.append(f"@@COMMAND pip install {pkg}\n")
                blocks.append(f"{sk.content.strip()}\n")

        if files:
            blocks.append("\n## EXISTING PROJECT WORKSPACE FILES\n")
            blocks.append(f"The project directory already contains {len(files)} file(s):\n")
            for f in files[:30]:
                blocks.append(f"- `{f['path']}` ({f['size']} bytes)")

            embedded_count = 0
            for f in files:
                if embedded_count >= 5:
                    break
                if f["size"] < 12_000 and f["ext"] in (
                    ".py",
                    ".js",
                    ".ts",
                    ".json",
                    ".html",
                    ".css",
                    ".md",
                    ".txt",
                ):
                    full_p = os.path.join(workspace_root, f["path"])
                    try:
                        with open(full_p, "r", encoding="utf-8") as fh:
                            content = fh.read()
                        blocks.append(f"\n### File: `{f['path']}`\n```\n{content}\n```\n")
                        embedded_count += 1
                    except OSError:
                        pass
        else:
            blocks.append("\n## WORKSPACE STATE\nThis is a clean/new project workspace folder.\n")

        return "\n".join(blocks)

    @staticmethod
    def build_verification_prompt(
        goal: str,
        plan: AgentPlan,
        result: ExecutionResult,
        workspace_root: str,
    ) -> str:
        """Builds the post-execution verification or diagnostic feedback prompt.

        Args:
            goal: Original user goal.
            plan: The AgentPlan that was executed.
            result: Subprocess execution outcome.
            workspace_root: Project directory.

        Returns:
            The feedback prompt for the LLM.
        """
        status_headline = (
            "✅ EXECUTION SUCCEEDED (Exit code 0)"
            if result.success
            else f"❌ EXECUTION FAILED (Exit code {result.exit_code})"
        )

        blocks: List[str] = [
            f"# SURGEON AGENT PROTOCOL v{AGENT_PROTOCOL_VERSION} — POST-EXECUTION FEEDBACK\n",
            f"Status: {status_headline}\n",
            f"Executed Command: `{result.command}` in {result.duration_sec:.2f}s\n",
            "\n## ORIGINAL GOAL\n",
            goal.strip(),
            "\n",
        ]

        if result.stdout:
            blocks.append("\n## STDOUT LOGS\n```text\n" + result.stdout[-4000:] + "\n```\n")
        if result.stderr:
            blocks.append("\n## STDERR / TRACEBACK LOGS\n```text\n" + result.stderr[-4000:] + "\n```\n")

        if result.artifacts_found:
            blocks.append("\n## CREATED ARTIFACTS / DELIVERABLES\n")
            for art in result.artifacts_found:
                blocks.append(f"- `{art}`")
            blocks.append("\n")

        if not result.success:
            blocks.append(
                "## REPAIR INSTRUCTIONS\n\n"
                "The code encountered an error or traceback as shown above.\n"
                "Please diagnose the issue and provide the fix using:\n"
                "1. An <EXPLANATION> of what caused the bug and how you fixed it.\n"
                "2. Any additional @@COMMAND setup lines if a dependency was missing.\n"
                "3. @@FILE or surgical @@EDIT blocks with the corrected code.\n"
                "4. A @@RUN line with the command to re-run the code.\n"
            )
        else:
            blocks.append(
                "## VERIFICATION TASK\n\n"
                "The command executed without errors and produced the outputs above.\n"
                "1. Confirm if the original user goal is fully accomplished.\n"
                "2. If any follow-up refinements, tests, or features are recommended, "
                "supply the @@FILE / @@EDIT / @@RUN blocks accordingly.\n"
            )

        return "\n".join(blocks)


# --------------------------------------------------------------------------- #
# Environment Secrets & Variables (.env)
# --------------------------------------------------------------------------- #


class EnvManager:
    """Manages project environment variables (.env) safely."""

    @staticmethod
    def parse_env(text: str) -> Dict[str, str]:
        """Parses a standard .env file content into a dictionary."""
        env_vars: Dict[str, str] = {}
        for raw_line in text.splitlines():
            line = raw_line.strip()
            if not line or line.startswith("#"):
                continue
            if "=" in line:
                key, val = line.split("=", 1)
                k = key.strip()
                v = val.strip()
                if (v.startswith('"') and v.endswith('"')) or (v.startswith("'") and v.endswith("'")):
                    v = v[1:-1]
                if k:
                    env_vars[k] = v
        return env_vars

    @staticmethod
    def format_env(env_vars: Dict[str, str]) -> str:
        """Formats a dictionary into standard .env file content."""
        lines = ["# Project Environment Variables (.env)", "# Generated by Text Surgeon Agent", ""]
        for k in sorted(env_vars.keys()):
            val = env_vars[k]
            if " " in val or "\n" in val or '"' in val:
                escaped = val.replace('"', '\\"')
                lines.append(f'{k}="{escaped}"')
            else:
                lines.append(f"{k}={val}")
        return "\n".join(lines) + "\n"

    @classmethod
    def load(cls, workspace_root: str) -> Dict[str, str]:
        """Loads .env file from workspace root."""
        env_path = os.path.join(workspace_root, ".env")
        if not os.path.isfile(env_path):
            return {}
        try:
            with open(env_path, "r", encoding="utf-8", errors="replace") as fh:
                return cls.parse_env(fh.read())
        except OSError:
            return {}

    @classmethod
    def save(cls, workspace_root: str, env_vars: Dict[str, str]) -> None:
        """Saves dictionary to .env file in workspace root."""
        env_path = os.path.join(workspace_root, ".env")
        content = cls.format_env(env_vars)
        with open(env_path, "w", encoding="utf-8", newline="\n") as fh:
            fh.write(content)


# --------------------------------------------------------------------------- #
# ZIP Project Exporter
# --------------------------------------------------------------------------- #


def export_project_zip(workspace_root: str) -> bytes:
    """Packs the project workspace into a clean ZIP archive (in-memory)."""
    clean_root = os.path.abspath(workspace_root)
    if not os.path.isdir(clean_root):
        raise AgentError(f"Directory '{workspace_root}' does not exist.")

    ignore_dirs = {
        ".surgeon",
        ".git",
        "__pycache__",
        "node_modules",
        ".venv",
        "venv",
        ".pytest_cache",
        ".idea",
        ".vscode",
    }

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(clean_root):
            dirs[:] = [d for d in dirs if d not in ignore_dirs]
            for fname in sorted(files):
                if fname.startswith(".surgeon") or fname.endswith(".tmp") or fname.endswith(".bak"):
                    continue
                full_path = os.path.join(root, fname)
                arcname = os.path.relpath(full_path, clean_root).replace("\\", "/")
                try:
                    zf.write(full_path, arcname)
                except OSError:
                    continue

    buf.seek(0)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# Nexus-Adopted Smart Key Manager & Multi-Key Rotation Engine
# --------------------------------------------------------------------------- #

# Error classification kinds (matching Nexus §4.2)
QUOTA = "quota"
BAD_KEY = "bad_key"
MODEL = "model"
TRANSIENT = "transient"
UNKNOWN = "unknown"

RPM_COOLDOWN_SEC: Final[int] = 75
_DAY_SECONDS: Final[int] = 86400

# Model quota metadata from Nexus registry
NEXUS_MODEL_CAPACITIES: Final[Dict[str, Dict[str, int]]] = {
    "gemini-2.5-flash": {"rpm": 15, "rpd": 1500},
    "gemini-2.5-flash-lite": {"rpm": 15, "rpd": 1500},
    "gemini-2.0-flash": {"rpm": 15, "rpd": 1500},
    "gemini-1.5-flash": {"rpm": 15, "rpd": 1500},
    "gemini-1.5-pro": {"rpm": 5, "rpd": 50},
    "gemini-2.5-pro": {"rpm": 5, "rpd": 50},
    "gemini-3.1-flash-lite": {"rpm": 15, "rpd": 500},
    "gemini-3.5-flash-lite": {"rpm": 15, "rpd": 500},
    "gemini-3.6-flash": {"rpm": 5, "rpd": 20},
    "gemini-3.5-flash": {"rpm": 5, "rpd": 20},
    "gemini-3-flash": {"rpm": 5, "rpd": 20},
}


def _today_utc() -> str:
    return time.strftime("%Y-%m-%d", time.gmtime())


def _next_utc_midnight() -> float:
    now = time.time()
    return (int(now // _DAY_SECONDS) + 1) * _DAY_SECONDS


def _kid(key: str) -> str:
    """Non-reversible stable hash id for key telemetry and storage."""
    return hashlib.sha256(key.encode("utf-8")).hexdigest()[:12]


@dataclass
class KeyStats:
    """Telemetry and rate-limit tracking for a single API key (Nexus-compatible)."""

    total_requests: int = 0
    successful_requests: int = 0
    failed_requests: int = 0
    cooldown_until: float = 0.0
    cooldown_reason: str = ""
    invalid: bool = False
    invalid_reason: str = ""
    last_kind: str = "ok"
    model_usage: Dict[str, int] = field(default_factory=dict)


class KeyManager:
    """Smart multi-key manager adopting the complete Nexus API rotation engine."""

    _stats: Dict[str, Dict[str, KeyStats]] = {}  # {provider: {kid: KeyStats}}
    _lock = threading.RLock()
    _day: str = _today_utc()

    @classmethod
    def _roll_day(cls) -> None:
        today = _today_utc()
        if today != cls._day:
            cls._day = today
            # Reset daily usage counts
            for prov in cls._stats.values():
                for st in prov.values():
                    st.model_usage.clear()

    @classmethod
    def _get_key_stats(cls, provider: str, key: str) -> KeyStats:
        clean_prov = provider.lower().strip()
        kid = _kid(key)
        cls._roll_day()
        if clean_prov not in cls._stats:
            cls._stats[clean_prov] = {}
        if kid not in cls._stats[clean_prov]:
            cls._stats[clean_prov][kid] = KeyStats()
        return cls._stats[clean_prov][kid]

    @classmethod
    def mask_key(cls, key: str) -> str:
        """Returns masked representation for safe UI rendering."""
        clean = key.strip()
        if len(clean) <= 8:
            return "••••"
        return f"••••{clean[-4:]}"

    @classmethod
    def parse_keys(cls, raw: Union[str, List[str], Sequence[str]]) -> List[str]:
        """Parses comma-separated, semicolon-separated, newline, or list of API keys."""
        if isinstance(raw, (list, tuple, set)):
            keys: List[str] = []
            for item in raw:
                keys.extend(cls.parse_keys(str(item)))
            return [k for k in keys if k]
        text = str(raw or "").strip()
        if not text:
            return []
        tokens = [k.strip() for k in re.split(r"[\r\n,;]+", text) if k.strip()]
        return tokens

    @classmethod
    def load_default_keys(cls, search_dirs: Optional[Sequence[str]] = None) -> List[str]:
        """Discovers and loads default API keys from apis.txt, apis.text, or .env files."""
        candidates: List[str] = []
        app_dir = os.path.dirname(os.path.abspath(__file__))
        check_dirs = [app_dir, os.getcwd()]
        if search_dirs:
            check_dirs.extend(search_dirs)

        for d in check_dirs:
            if not d or not os.path.isdir(d):
                continue
            for fname in ("apis.txt", "apis.text", "api_keys.txt", ".env"):
                fpath = os.path.join(d, fname)
                if os.path.isfile(fpath):
                    try:
                        with open(fpath, "r", encoding="utf-8", errors="replace") as fh:
                            content = fh.read()
                        if fname.startswith(".env"):
                            env_map = EnvManager.parse_env(content)
                            for k in ("GEMINI_API_KEYS", "GEMINI_API_KEY", "OPENAI_API_KEY", "API_KEYS"):
                                if k in env_map:
                                    candidates.extend(cls.parse_keys(env_map[k]))
                        else:
                            candidates.extend(cls.parse_keys(content))
                    except OSError:
                        pass

        # Also check environment variables
        for env_var in ("GEMINI_API_KEYS", "GEMINI_API_KEY", "OPENAI_API_KEY", "API_KEYS"):
            val = os.environ.get(env_var)
            if val:
                candidates.extend(cls.parse_keys(val))

        seen = set()
        deduped: List[str] = []
        for k in candidates:
            if k not in seen:
                seen.add(k)
                deduped.append(k)
        return deduped

    @classmethod
    def get_cap_for(cls, model: str) -> Optional[int]:
        """Returns daily cap (RPD) for a model."""
        clean_m = model.lower().strip()
        row = NEXUS_MODEL_CAPACITIES.get(clean_m)
        return row.get("rpd") if row else None

    @classmethod
    def classify_error(cls, status_code: int, error_body: str) -> str:
        """Maps HTTP status code / SDK error to Nexus error kinds."""
        s = f"status {status_code} {error_body}".lower()
        if any(x in s for x in ("resource_exhausted", "429", "quota", "rate limit", "rate-limit", "exhausted")):
            return QUOTA
        if any(x in s for x in ("api_key_invalid", "api key not valid", "invalid api key", "401", "unauthenticated", "permission_denied", "403", "invalid authentication")):
            return BAD_KEY
        if any(x in s for x in ("not_found", "not found", "404", "is not supported", "unsupported", "does not exist")):
            return MODEL
        if any(x in s for x in ("500", "502", "503", "504", "internal", "unavailable", "deadline", "timeout", "timed out", "connection", "temporarily")):
            return TRANSIENT
        return UNKNOWN

    @classmethod
    def order_keys(cls, provider: str, keys: Sequence[str], model: str = "") -> List[str]:
        """Orders keys Nexus-style: valid first, not in cooldown, most remaining capacity first."""
        clean_keys = [k.strip() for k in keys if k and k.strip()]
        if not clean_keys:
            return []

        now = time.time()
        cap = cls.get_cap_for(model)

        with cls._lock:
            cls._roll_day()

            def rank(k: str) -> Tuple[bool, bool, int]:
                st = cls._get_key_stats(provider, k)
                used = st.model_usage.get(model, 0)
                in_cd = st.cooldown_until > now
                remaining = (cap - used) if cap else 10_000
                return (st.invalid, in_cd, -remaining)

            return sorted(clean_keys, key=rank)

    @classmethod
    def get_active_key(cls, keys: Sequence[str], provider: str = "openai") -> Optional[str]:
        """Returns best active key from pool according to health and capacity."""
        ordered = cls.order_keys(provider, keys)
        return ordered[0] if ordered else None

    @classmethod
    def record_success(cls, provider: str, key: str, model: str = "") -> None:
        """Records successful response on key."""
        with cls._lock:
            cls._roll_day()
            st = cls._get_key_stats(provider, key)
            st.total_requests += 1
            st.successful_requests += 1
            st.last_kind = "ok"
            if model:
                st.model_usage[model] = st.model_usage.get(model, 0) + 1

    @classmethod
    def record_error(
        cls,
        provider: str,
        key: str,
        model: str,
        status_code: int,
        error_body: str,
        cooldown_sec: int = RPM_COOLDOWN_SEC,
    ) -> str:
        """Records error on key, applies appropriate cooldown / invalidation, returns error kind."""
        with cls._lock:
            cls._roll_day()
            st = cls._get_key_stats(provider, key)
            st.total_requests += 1
            st.failed_requests += 1
            kind = cls.classify_error(status_code, error_body)
            st.last_kind = kind

            if kind == BAD_KEY:
                st.invalid = True
                st.invalid_reason = f"HTTP {status_code}: {error_body[:200]}"
            elif kind == QUOTA:
                cap = cls.get_cap_for(model)
                used = st.model_usage.get(model, 0)
                if cap and used >= cap:
                    st.cooldown_until = _next_utc_midnight()
                    st.cooldown_reason = f"Daily cap reached ({used}/{cap}). Cooldown until UTC midnight."
                else:
                    st.cooldown_until = time.time() + cooldown_sec
                    st.cooldown_reason = f"RPM limit exceeded. Cooldown {cooldown_sec}s."
            return kind

    @classmethod
    def record_key_status(
        cls,
        key: str,
        status_code: int,
        provider: str = "openai",
        model: str = "",
        cooldown_sec: int = RPM_COOLDOWN_SEC,
    ) -> str:
        return cls.record_error(
            provider=provider,
            key=key,
            model=model,
            status_code=status_code,
            error_body=f"HTTP {status_code}",
            cooldown_sec=cooldown_sec,
        )

    @classmethod
    def get_status(cls, provider: str, keys: Sequence[str]) -> List[Dict[str, Any]]:
        """Returns live status of keys for UI monitoring."""
        res: List[Dict[str, Any]] = []
        now = time.time()
        with cls._lock:
            cls._roll_day()
            for k in keys:
                if not k:
                    continue
                st = cls._get_key_stats(provider, k)
                state = "ready"
                if st.invalid:
                    state = "invalid"
                elif st.cooldown_until > now:
                    state = "cooldown"
                res.append(
                    {
                        "masked": cls.mask_key(k),
                        "state": state,
                        "total_requests": st.total_requests,
                        "successful_requests": st.successful_requests,
                        "failed_requests": st.failed_requests,
                        "cooldown_remaining": max(0, int(st.cooldown_until - now)),
                        "invalid_reason": st.invalid_reason,
                        "last_kind": st.last_kind,
                        "model_usage": dict(st.model_usage),
                    }
                )
        return res

    @classmethod
    def get_pool_status(cls, keys: Sequence[str], provider: str = "openai") -> Dict[str, Any]:
        """Returns summary counts of keys."""
        status_list = cls.get_status(provider, keys)
        ready_count = sum(1 for s in status_list if s["state"] == "ready")
        cooldown_count = sum(1 for s in status_list if s["state"] == "cooldown")
        invalid_count = sum(1 for s in status_list if s["state"] == "invalid")
        return {
            "total": len(status_list),
            "ready": ready_count,
            "cooldown": cooldown_count,
            "invalid": invalid_count,
            "keys": status_list,
        }


# --------------------------------------------------------------------------- #
# Pure Standard-Library LLM Client & Autonomous Auto-Pilot
# --------------------------------------------------------------------------- #


@dataclass
class LLMConfig:
    """Configuration for LLM API connection."""

    provider: str = "gemini"  # gemini | openai | groq | deepseek | anthropic | openrouter | ollama | custom
    model: str = "gemini-3.1-flash-lite"
    api_key: str = ""
    api_keys: List[str] = field(default_factory=list)
    base_url: str = "http://localhost:11434"
    temperature: float = 0.2
    max_tokens: int = 4096
    timeout_sec: int = 120

    def get_all_keys(self) -> List[str]:
        """Returns all configured API keys."""
        keys: List[str] = []
        if self.api_keys:
            keys.extend(KeyManager.parse_keys(self.api_keys))
        if self.api_key:
            keys.extend(KeyManager.parse_keys(self.api_key))
        if not keys:
            keys.extend(KeyManager.load_default_keys())
        seen = set()
        deduped = []
        for k in keys:
            if k not in seen:
                seen.add(k)
                deduped.append(k)
        return deduped


class LLMClient:
    """Nexus-style LLM Client with cached GenAI client reuse, REST fallbacks, and smart rotation."""

    _clients: Dict[str, Any] = {}
    _client_lock = threading.Lock()

    @classmethod
    def _get_genai_client(cls, api_key: str) -> Optional[Any]:
        """Reuses cached genai.Client per key (fixes SDK transport closure bug)."""
        with cls._client_lock:
            if api_key in cls._clients:
                return cls._clients[api_key]
            try:
                from google import genai
                c = genai.Client(api_key=api_key)
                cls._clients[api_key] = c
                return c
            except Exception:
                return None

    @classmethod
    def list_ollama_models(cls, base_url: str = "http://localhost:11434", timeout_sec: int = 4) -> List[str]:
        """Lists available local models running in Ollama."""
        clean_url = base_url.rstrip("/") + "/api/tags"
        req = urllib.request.Request(clean_url, headers={"User-Agent": "TextSurgeon/2.0"})
        try:
            with urllib.request.urlopen(req, timeout=timeout_sec) as resp:
                if resp.status == 200:
                    data = json.loads(resp.read().decode("utf-8"))
                    return [m.get("name", "") for m in data.get("models", []) if m.get("name")]
        except Exception:
            return []
        return []

    @classmethod
    def send_prompt(
        cls,
        prompt: str,
        config: LLMConfig,
        system_prompt: Optional[str] = None,
    ) -> str:
        """Sends a prompt using Nexus multi-key rotation, failover, and model error detection."""
        provider = config.provider.lower().strip()

        if provider == "ollama":
            return cls._call_ollama(prompt, config, system_prompt)

        all_keys = config.get_all_keys()
        if not all_keys:
            if provider == "custom":
                return cls._call_openai_compatible(prompt, config, system_prompt, key="")
            raise AgentError(
                f"No API key provided for {provider.title()}.",
                hint=f"Please enter one or more API keys for {provider.title()}.",
            )

        ordered_keys = KeyManager.order_keys(provider, all_keys, model=config.model)
        if not ordered_keys:
            raise AgentError(
                f"All configured API keys for {provider.title()} are invalid or in cooldown.",
                hint="Please check your API keys or wait for cooldown to expire.",
            )

        last_error: Optional[Exception] = None

        for key in ordered_keys:
            # Nexus attempt loop: 1 quick retry for transient blips before rotating
            for attempt in (1, 2):
                try:
                    if provider == "gemini":
                        reply = cls._call_gemini(prompt, config, system_prompt, key=key)
                    elif provider == "anthropic":
                        reply = cls._call_anthropic(prompt, config, system_prompt, key=key)
                    else:
                        reply = cls._call_openai_compatible(prompt, config, system_prompt, key=key)

                    KeyManager.record_success(provider, key, model=config.model)
                    return reply

                except urllib.error.HTTPError as exc:
                    err_body = exc.read().decode("utf-8", errors="replace")
                    kind = KeyManager.record_error(
                        provider=provider,
                        key=key,
                        model=config.model,
                        status_code=exc.code,
                        error_body=err_body,
                    )
                    last_error = AgentError(
                        f"{provider.title()} API Error ({exc.code}): {err_body[:250]}",
                        hint=f"Key {KeyManager.mask_key(key)} recorded error kind '{kind}'.",
                    )
                    # If model not found (404), stop early — all keys would fail the same way
                    if kind == MODEL:
                        raise AgentError(f"Model '{config.model}' was rejected by {provider.title()}: {err_body[:250]}") from exc
                    if kind == TRANSIENT and attempt == 1:
                        time.sleep(1.5)
                        continue
                    break  # rotate to next key

                except urllib.error.URLError as exc:
                    last_error = AgentError(f"Network connection error to {provider.title()}: {exc}")
                    if attempt == 1:
                        time.sleep(1.5)
                        continue
                    break

                except Exception as exc:
                    last_error = exc
                    if attempt == 1:
                        time.sleep(1.0)
                        continue
                    break

        if last_error:
            if isinstance(last_error, AgentError):
                raise last_error
            raise AgentError(f"{provider.title()} call failed across all keys: {last_error}") from last_error

        raise AgentError(f"Failed to receive response from {provider.title()}.")

    @classmethod
    def _call_ollama(cls, prompt: str, config: LLMConfig, system_prompt: Optional[str]) -> str:
        url = config.base_url.rstrip("/") + "/api/generate"
        payload: Dict[str, Any] = {
            "model": config.model or "llama3",
            "prompt": prompt,
            "stream": False,
            "options": {
                "temperature": config.temperature,
                "num_predict": config.max_tokens,
            },
        }
        if system_prompt:
            payload["system"] = system_prompt

        req_body = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request(
            url,
            data=req_body,
            headers={"Content-Type": "application/json", "User-Agent": "TextSurgeon/2.0"},
            method="POST",
        )
        try:
            with urllib.request.urlopen(req, timeout=config.timeout_sec) as resp:
                data = json.loads(resp.read().decode("utf-8"))
                return str(data.get("response", "")).strip()
        except urllib.error.URLError as exc:
            raise AgentError(
                f"Failed to connect to Ollama at {url}: {exc}",
                hint="Ensure Ollama is running (`ollama serve`) and the model is pulled (`ollama pull llama3`).",
            ) from exc

    @classmethod
    def _call_gemini(
        cls,
        prompt: str,
        config: LLMConfig,
        system_prompt: Optional[str],
        key: str,
    ) -> str:
        model = config.model or "gemini-3.1-flash-lite"

        # 1. Try google-genai SDK if available
        client = cls._get_genai_client(key)
        if client is not None:
            try:
                from google.genai import types
                gen_cfg = types.GenerateContentConfig(
                    temperature=config.temperature,
                    max_output_tokens=config.max_tokens,
                )
                if system_prompt:
                    gen_cfg.system_instruction = system_prompt
                resp = client.models.generate_content(model=model, contents=prompt, config=gen_cfg)
                if resp and resp.text:
                    return resp.text.strip()
            except Exception as e:
                # If SDK encounters non-HTTP Python error, fall through to REST call
                pass

        # 2. REST Endpoint fallback (Zero dependency standard library)
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={key}"
        parts = []
        if system_prompt:
            parts.append({"text": system_prompt + "\n\n" + prompt})
        else:
            parts.append({"text": prompt})

        payload = {
            "contents": [{"parts": parts}],
            "generationConfig": {
                "temperature": config.temperature,
                "maxOutputTokens": config.max_tokens,
            },
        }
        req_body = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request(
            url,
            data=req_body,
            headers={"Content-Type": "application/json", "User-Agent": "TextSurgeon/2.0"},
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=config.timeout_sec) as resp:
            data = json.loads(resp.read().decode("utf-8"))
            candidates = data.get("candidates", [])
            if candidates:
                parts_out = candidates[0].get("content", {}).get("parts", [])
                if parts_out:
                    return str(parts_out[0].get("text", "")).strip()
            raise AgentError("No content returned from Gemini API.")

    @classmethod
    def _call_openai_compatible(
        cls,
        prompt: str,
        config: LLMConfig,
        system_prompt: Optional[str],
        key: str,
    ) -> str:
        base_urls = {
            "openai": "https://api.openai.com/v1/chat/completions",
            "groq": "https://api.groq.com/openai/v1/chat/completions",
            "deepseek": "https://api.deepseek.com/v1/chat/completions",
            "openrouter": "https://openrouter.ai/api/v1/chat/completions",
        }
        url = config.base_url if config.provider == "custom" else base_urls.get(config.provider, "https://api.openai.com/v1/chat/completions")
        if not url.endswith("/chat/completions") and config.provider == "custom":
            url = url.rstrip("/") + "/chat/completions"

        messages: List[Dict[str, str]] = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": prompt})

        default_model = "gpt-4o"
        if config.provider == "groq":
            default_model = "llama-3.3-70b-versatile"
        elif config.provider == "deepseek":
            default_model = "deepseek-chat"

        payload = {
            "model": config.model or default_model,
            "messages": messages,
            "temperature": config.temperature,
            "max_tokens": config.max_tokens,
        }

        headers = {
            "Content-Type": "application/json",
            "User-Agent": "TextSurgeon/2.0",
        }
        if key:
            headers["Authorization"] = f"Bearer {key}"
        if config.provider == "openrouter":
            headers["HTTP-Referer"] = "https://github.com/shady/text-surgeon"
            headers["X-Title"] = "Text Surgeon Agent"

        req_body = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request(url, data=req_body, headers=headers, method="POST")
        with urllib.request.urlopen(req, timeout=config.timeout_sec) as resp:
            data = json.loads(resp.read().decode("utf-8"))
            choices = data.get("choices", [])
            if choices and "message" in choices[0]:
                return str(choices[0]["message"].get("content", "")).strip()
            raise AgentError("No response choices returned by OpenAI API.")

    @classmethod
    def _call_anthropic(
        cls,
        prompt: str,
        config: LLMConfig,
        system_prompt: Optional[str],
        key: str,
    ) -> str:
        url = "https://api.anthropic.com/v1/messages"
        payload: Dict[str, Any] = {
            "model": config.model or "claude-3-5-sonnet-20241022",
            "max_tokens": config.max_tokens,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": config.temperature,
        }
        if system_prompt:
            payload["system"] = system_prompt

        headers = {
            "Content-Type": "application/json",
            "x-api-key": key,
            "anthropic-version": "2023-06-01",
            "User-Agent": "TextSurgeon/2.0",
        }
        req_body = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request(url, data=req_body, headers=headers, method="POST")
        with urllib.request.urlopen(req, timeout=config.timeout_sec) as resp:
            data = json.loads(resp.read().decode("utf-8"))
            content_blocks = data.get("content", [])
            if content_blocks and "text" in content_blocks[0]:
                return str(content_blocks[0]["text"]).strip()
            raise AgentError("No content blocks returned from Anthropic API.")


# Common Python module-to-pip mapping for autonomous dependency healing
_MODULE_TO_PIP_MAP: Final[Dict[str, str]] = {
    "pptx": "python-pptx",
    "docx": "python-docx",
    "bs4": "beautifulsoup4",
    "dotenv": "python-dotenv",
    "PIL": "pillow",
    "cv2": "opencv-python",
    "yaml": "pyyaml",
    "sklearn": "scikit-learn",
    "telebot": "pyTelegramBotAPI",
    "telegram": "python-telegram-bot",
    "pandas": "pandas",
    "openpyxl": "openpyxl",
    "matplotlib": "matplotlib",
    "seaborn": "seaborn",
    "numpy": "numpy",
    "requests": "requests",
    "aiohttp": "aiohttp",
    "fastapi": "fastapi",
    "uvicorn": "uvicorn",
}


class AutoPilotEngine:
    """Executes multi-round autonomous agent loops with self-healing and smart dependency resolution."""

    @classmethod
    def run_autonomous_loop(
        cls,
        workspace_root: str,
        goal: str,
        llm_config: LLMConfig,
        runtime: str = "python",
        max_rounds: int = 3,
        include_memory: bool = True,
        skill_names: Optional[Sequence[str]] = None,
        auto_detect_skills: bool = True,
        on_round_callback: Optional[Callable[[Dict[str, Any]], None]] = None,
    ) -> Dict[str, Any]:
        """Runs an autonomous multi-round loop until exit 0 or max_rounds reached."""
        ws = AgentWorkspace(workspace_root)
        rounds_history: List[Dict[str, Any]] = []
        current_prompt = AgentPromptBuilder.build_task_prompt(
            goal=goal,
            workspace_root=workspace_root,
            runtime=runtime,
            include_memory=include_memory,
            skill_names=skill_names,
            auto_detect_skills=auto_detect_skills,
        )
        last_result: Optional[ExecutionResult] = None

        for round_idx in range(1, max_rounds + 1):
            round_start = time.perf_counter()
            round_info: Dict[str, Any] = {
                "round": round_idx,
                "prompt": current_prompt,
                "timestamp": datetime.now().isoformat(),
            }

            try:
                # 1. Query LLM using Nexus API Caller
                raw_reply = LLMClient.send_prompt(current_prompt, llm_config)
                round_info["raw_reply"] = raw_reply

                # 2. Parse response
                plan = parse_agent_response(raw_reply)
                round_info["plan"] = {
                    "explanation": plan.explanation,
                    "files_count": len(plan.file_actions),
                    "setup_commands": plan.setup_commands,
                    "run_command": plan.run_command,
                    "expected_artifacts": plan.expected_artifacts,
                }
                round_info["explanation"] = plan.explanation
                round_info["files"] = [fa.path for fa in plan.file_actions]

                # 3. Apply changes with snapshot backup
                apply_res = ws.apply_plan(plan, backup=True)
                round_info["apply"] = apply_res

                # 4. Run setup commands if any (e.g. pip install)
                setup_logs: List[str] = []
                for scmd in plan.setup_commands:
                    s_res = ws.execute_command(scmd, timeout=180)
                    setup_logs.append(f"$ {scmd}\n{s_res.stdout}\n{s_res.stderr}")
                round_info["setup_logs"] = setup_logs

                # 5. Run main command
                cmd_to_run = plan.run_command or ("python main.py" if runtime == "python" else "node index.js")
                exec_res = ws.execute_command(cmd_to_run, expected_artifacts=plan.expected_artifacts)

                # 6. Autonomous Self-Healing for Missing Python Modules
                if not exec_res.success and runtime == "python":
                    missing_mod_match = re.search(r"No module named ['\"]([a-zA-Z0-9_\-]+)['\"]", exec_res.stderr)
                    if missing_mod_match:
                        raw_mod = missing_mod_match.group(1)
                        pkg_to_install = _MODULE_TO_PIP_MAP.get(raw_mod, raw_mod)
                        install_res = ws.install_dependencies([pkg_to_install])
                        round_info["auto_installed_package"] = pkg_to_install
                        round_info["auto_install_log"] = f"$ pip install {pkg_to_install}\n{install_res.stdout}\n{install_res.stderr}"
                        # Re-run after package install
                        exec_res = ws.execute_command(cmd_to_run, expected_artifacts=plan.expected_artifacts)

                last_result = exec_res
                round_info["execution"] = {
                    "command": exec_res.command,
                    "exit_code": exec_res.exit_code,
                    "stdout": exec_res.stdout,
                    "stderr": exec_res.stderr,
                    "duration_sec": exec_res.duration_sec,
                    "artifacts_found": exec_res.artifacts_found,
                    "success": exec_res.success,
                }

                round_info["duration_sec"] = time.perf_counter() - round_start
                rounds_history.append(round_info)
                ws.record_round(round_info)

                if on_round_callback:
                    on_round_callback(round_info)

                # If successful, exit loop!
                if exec_res.success:
                    ws.memory.record_short_term_event(
                        event_type="goal_completed",
                        goal=goal,
                        details={
                            "rounds_completed": round_idx,
                            "artifacts": exec_res.artifacts_found,
                            "command": exec_res.command,
                            "duration_sec": round_info.get("duration_sec", 0.0),
                        },
                    )
                    return {
                        "status": "success",
                        "success": True,
                        "rounds_completed": round_idx,
                        "final_round": round_info,
                        "rounds": rounds_history,
                        "artifacts": exec_res.artifacts_found,
                        "final_artifacts": exec_res.artifacts_found,
                    }

                # Build verification / diagnostic repair prompt for next round
                current_prompt = AgentPromptBuilder.build_verification_prompt(
                    goal=goal,
                    plan=plan,
                    result=exec_res,
                    workspace_root=workspace_root,
                )

            except Exception as exc:
                round_info["error"] = str(exc)
                round_info["duration_sec"] = time.perf_counter() - round_start
                rounds_history.append(round_info)
                ws.record_round(round_info)
                ws.memory.record_short_term_event(
                    event_type="auto_pilot_error",
                    goal=goal,
                    details={"error": str(exc), "round": round_idx},
                )
                if on_round_callback:
                    on_round_callback(round_info)

                # Retry next round with error notice
                current_prompt = (
                    f"# SURGEON AGENT PROTOCOL v{AGENT_PROTOCOL_VERSION} — ERROR IN LAST ATTEMPT\n\n"
                    f"An error occurred while applying your plan:\n```text\n{exc}\n```\n\n"
                    "Please re-generate the complete solution adhering to @@FILE and @@RUN blocks.\n\n"
                    f"## GOAL\n{goal}\n"
                )

        return {
            "status": "max_rounds_reached",
            "success": False,
            "rounds_completed": max_rounds,
            "final_round": rounds_history[-1] if rounds_history else {},
            "rounds": rounds_history,
            "artifacts": last_result.artifacts_found if last_result else [],
            "final_artifacts": last_result.artifacts_found if last_result else [],
        }

